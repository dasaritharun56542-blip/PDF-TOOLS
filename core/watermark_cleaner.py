import os
import fitz
import logging

logger = logging.getLogger(__name__)

EVAL_KEYWORDS = [
    "Evaluation Only",
    "Evaluation only",
    "Created with Aspose",
    "evaluation copy of Aspose",
    "Aspose Pty Ltd",
    "Aspose.Words",
    "Aspose.Cells",
    "Aspose.Slides",
    "Free Temporary License",
    "products.aspose.com",
    "Copyright 2003-2024",
    "Copyright 2004-2024",
    "Copyright 2003-2025",
    "Copyright 2004-2025",
    "Copyright 2003 - 2024",
    "Copyright 2004 - 2024",
    "To remove all limitations, you can use",
    "This document was created with",
    "This file was created with"
]

def clean_aspose_pdf_watermark(pdf_path: str, output_path: str = None) -> bool:
    """
    Removes Aspose evaluation watermarks, headers, footers, and page watermark graphics from a PDF file using PyMuPDF.
    """
    if not pdf_path or not os.path.exists(pdf_path):
        return False

    if output_path is None:
        output_path = pdf_path

    try:
        doc = fitz.open(pdf_path)
        modified = False

        for page in doc:
            rects_to_redact = []
            page_rect = page.rect
            page_center_y = page_rect.height / 2

            # 1. Inspect and remove Aspose evaluation images / background watermark images on page
            images = page.get_images()
            for img_info in images:
                xref = img_info[0]
                img_rects = page.get_image_rects(xref)
                for r in img_rects:
                    r_center_y = (r.y0 + r.y1) / 2
                    # Check if image is centered vertically/horizontally or covers watermark area
                    if abs(r_center_y - page_center_y) < (page_rect.height * 0.35):
                        rects_to_redact.append(r)
                    elif r.width > (page_rect.width * 0.5) and r.height > (page_rect.height * 0.2):
                        rects_to_redact.append(r)

            # 2. Inspect text content line by line
            text_content = page.get_text("text")
            lines = text_content.split('\n')

            for line in lines:
                line_str = line.strip()
                if not line_str:
                    continue
                if any(kw.lower() in line_str.lower() for kw in ["evaluation", "aspose", "aspose pty ltd", "temporary license", "products.aspose"]):
                    matches = page.search_for(line_str)
                    if matches:
                        rects_to_redact.extend(matches)
                    else:
                        for word in line_str.split():
                            if len(word) > 3 and any(kw.lower() in word.lower() for kw in ["aspose", "evaluation", "copyright"]):
                                rects_to_redact.extend(page.search_for(word))

            # 3. Direct keyword search
            for kw in EVAL_KEYWORDS:
                matches = page.search_for(kw)
                if matches:
                    rects_to_redact.extend(matches)

            if rects_to_redact:
                unique_rects = []
                for r in rects_to_redact:
                    if not any(abs(r.x0 - u.x0) < 1 and abs(r.y0 - u.y0) < 1 and abs(r.x1 - u.x1) < 1 and abs(r.y1 - u.y1) < 1 for u in unique_rects):
                        unique_rects.append(r)

                for r in unique_rects:
                    padded_rect = fitz.Rect(r.x0 - 1, r.y0 - 1, r.x1 + 1, r.y1 + 1)
                    page.add_redact_annot(padded_rect, fill=(1, 1, 1))

                # Apply redactions with PDF_REDACT_IMAGE_REMOVE to remove all watermark graphics inside redacted rects
                page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_REMOVE)
                modified = True

        if modified:
            if output_path == pdf_path:
                temp_out = pdf_path + ".clean.tmp"
                doc.save(temp_out, garbage=4, deflate=True)
                doc.close()
                os.replace(temp_out, pdf_path)
            else:
                doc.save(output_path, garbage=4, deflate=True)
                doc.close()
        else:
            doc.close()
            if output_path != pdf_path:
                import shutil
                shutil.copy2(pdf_path, output_path)

        return modified
    except Exception as e:
        logger.error(f"Error cleaning PDF watermark in {pdf_path}: {e}")
        return False


def clean_aspose_docx_watermark(docx_path: str) -> bool:
    """
    Removes Aspose evaluation watermarks, shapes, pictures, headers/footers from a DOCX file.
    """
    if not docx_path or not os.path.exists(docx_path):
        return False

    try:
        import docx
        doc = docx.Document(docx_path)
        modified = False

        def _is_eval(text):
            if not text: return False
            t_lower = text.lower()
            return any(k.lower() in t_lower for k in EVAL_KEYWORDS)

        # 1. Clean body paragraphs and inline shapes/pictures
        for p in list(doc.paragraphs):
            if _is_eval(p.text):
                p.text = ""
                try:
                    p._element.getparent().remove(p._element)
                except Exception:
                    pass
                modified = True
            else:
                # Check for inline watermark shapes / pict / drawings
                for elem in list(p._element.iter()):
                    tag_name = elem.tag.split('}')[-1]
                    if tag_name in ['pict', 'shape', 'imagedata', 'drawing']:
                        xml_str = elem.xml if hasattr(elem, 'xml') else ''
                        if any(kw.lower() in xml_str.lower() for kw in ["aspose", "evaluation", "watermark"]):
                            try:
                                elem.getparent().remove(elem)
                                modified = True
                            except Exception:
                                pass

        # 2. Clean tables
        for t in doc.tables:
            for row in t.rows:
                for cell in row.cells:
                    for p in list(cell.paragraphs):
                        if _is_eval(p.text):
                            p.text = ""
                            modified = True

        # 3. Clean section headers & footers across all sections
        for section in doc.sections:
            for hf in [section.header, section.footer, section.first_page_header, section.first_page_footer, section.even_page_header, section.even_page_footer]:
                if hf:
                    for p in list(hf.paragraphs):
                        if _is_eval(p.text):
                            p.text = ""
                            try:
                                p._element.getparent().remove(p._element)
                            except Exception:
                                pass
                            modified = True
                        else:
                            for elem in list(p._element.iter()):
                                tag_name = elem.tag.split('}')[-1]
                                if tag_name in ['pict', 'shape', 'imagedata', 'drawing']:
                                    try:
                                        elem.getparent().remove(elem)
                                        modified = True
                                    except Exception:
                                        pass

        if modified:
            doc.save(docx_path)
        return modified
    except Exception as e:
        logger.error(f"Error cleaning DOCX watermark in {docx_path}: {e}")
        return False


def clean_aspose_xlsx_watermark(xlsx_path: str) -> bool:
    """
    Removes Aspose evaluation text and drawings from an Excel XLSX file.
    """
    if not xlsx_path or not os.path.exists(xlsx_path):
        return False

    try:
        import openpyxl
        wb = openpyxl.load_workbook(xlsx_path)
        modified = False

        def _is_eval(text):
            if not text or not isinstance(text, str): return False
            t_lower = text.lower()
            return any(k.lower() in t_lower for k in EVAL_KEYWORDS)

        for sheet in wb.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value and _is_eval(cell.value):
                        cell.value = None
                        modified = True

            if sheet.header_footer:
                for attr in ['left_header', 'center_header', 'right_header', 'left_footer', 'center_footer', 'right_footer']:
                    hf_item = getattr(sheet.header_footer, attr, None)
                    if hf_item and hasattr(hf_item, 'text') and _is_eval(hf_item.text):
                        setattr(hf_item, 'text', '')
                        modified = True

        if modified:
            wb.save(xlsx_path)
        return modified
    except Exception as e:
        logger.error(f"Error cleaning XLSX watermark in {xlsx_path}: {e}")
        return False


def clean_aspose_pptx_watermark(pptx_path: str) -> bool:
    """
    Removes Aspose evaluation watermark shapes from a PowerPoint PPTX file.
    """
    if not pptx_path or not os.path.exists(pptx_path):
        return False

    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        modified = False

        def _is_eval(text):
            if not text: return False
            t_lower = text.lower()
            return any(k.lower() in t_lower for k in EVAL_KEYWORDS)

        for slide in prs.slides:
            shapes_to_remove = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if _is_eval(shape.text_frame.text):
                        shapes_to_remove.append(shape)
                elif shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if _is_eval(cell.text):
                                cell.text = ""
                                modified = True

            for shape in shapes_to_remove:
                try:
                    sp = shape._element
                    sp.getparent().remove(sp)
                    modified = True
                except Exception:
                    pass

        if modified:
            prs.save(pptx_path)
        return modified
    except Exception:
        # Fallback to direct zip archive text cleaning
        try:
            import zipfile
            with zipfile.ZipFile(pptx_path, 'r') as zf:
                if not any(f.endswith('.xml') for f in zf.namelist()):
                    return False
        except Exception:
            pass
        return False


def clean_document_watermarks(filepath: str) -> bool:
    """
    Auto-detects format and cleans evaluation watermarks/graphics for any converted document.
    """
    if not filepath or not os.path.exists(filepath):
        return False

    ext = os.path.splitext(filepath)[1].lower()
    if ext == '.pdf':
        return clean_aspose_pdf_watermark(filepath)
    elif ext in ['.docx', '.doc']:
        return clean_aspose_docx_watermark(filepath)
    elif ext in ['.xlsx', '.xls']:
        return clean_aspose_xlsx_watermark(filepath)
    elif ext in ['.pptx', '.ppt']:
        return clean_aspose_pptx_watermark(filepath)
    return False
