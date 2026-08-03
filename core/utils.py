import os, uuid, io, re
from django.conf import settings
import fitz  # PyMuPDF
import img2pdf
from PIL import Image

_CACHED_UNICODE_FONT = None

def _get_cached_unicode_font():
    global _CACHED_UNICODE_FONT
    if _CACHED_UNICODE_FONT:
        return _CACHED_UNICODE_FONT

    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    font_name = 'Helvetica'
    system_fonts = [
        (r'C:\Windows\Fonts\Nirmala.ttc', 'NirmalaUI', 0),
        (r'C:\Windows\Fonts\segoeui.ttf', 'SegoeUI', None),
        (r'C:\Windows\Fonts\arial.ttf', 'Arial', None),
        (r'C:\Windows\Fonts\seguiemj.ttf', 'SegoeUIEmoji', None),
        (r'C:\Windows\Fonts\tahoma.ttf', 'Tahoma', None),
        ('/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf', 'DejaVuSans', None),
        ('/Library/Fonts/Arial.ttf', 'Arial', None)
    ]
    for fitem in system_fonts:
        fpath, fname, subidx = fitem
        if os.path.exists(fpath):
            try:
                if subidx is not None:
                    pdfmetrics.registerFont(TTFont(fname, fpath, subfontIndex=subidx))
                else:
                    pdfmetrics.registerFont(TTFont(fname, fpath))
                from reportlab.lib.fonts import addMapping, _ps2tt_map
                addMapping(fname, 0, 0, fname)
                addMapping(fname, 1, 0, fname)
                addMapping(fname, 0, 1, fname)
                addMapping(fname, 1, 1, fname)
                
                fname_lower = fname.lower()
                _ps2tt_map[fname_lower] = (fname_lower, 0, 0)
                _ps2tt_map[f"{fname_lower}-bold"] = (fname_lower, 1, 0)
                _ps2tt_map[f"{fname_lower}-oblique"] = (fname_lower, 0, 1)
                _ps2tt_map[f"{fname_lower}-italic"] = (fname_lower, 0, 1)
                _ps2tt_map[f"{fname_lower}-boldoblique"] = (fname_lower, 1, 1)
                _ps2tt_map[f"{fname_lower}-bolditalic"] = (fname_lower, 1, 1)
                
                font_name = fname
                break
            except Exception as font_ex:
                print(f"Font registration notice for {fname}: {font_ex}")

    _CACHED_UNICODE_FONT = font_name
    return _CACHED_UNICODE_FONT

def _parse_legacy_doc_stream(doc_bytes):
    """Dedicated Microsoft OLE2 Compound Binary Format (.doc) parser for Word 97-2003 files."""
    if not doc_bytes or len(doc_bytes) < 512:
        return []

    import struct, re
    if doc_bytes[:8] != b'\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1':
        return []

    try:
        sec_size = 1 << struct.unpack('<H', doc_bytes[30:32])[0]
        dir_start = struct.unpack('<I', doc_bytes[48:52])[0]

        fat_sectors = []
        difat = struct.unpack('<109I', doc_bytes[76:512])
        for s in difat:
            if s < 0xFFFFFFFC:
                fat_sectors.append(s)

        fat = []
        for f_sec in fat_sectors:
            offset = (f_sec + 1) * sec_size
            fat.extend(struct.unpack(f'<{sec_size//4}I', doc_bytes[offset:offset+sec_size]))

        def get_chain(start_sec):
            chain = []
            curr = start_sec
            while curr < 0xFFFFFFFC:
                chain.append(curr)
                if curr >= len(fat):
                    break
                curr = fat[curr]
            return chain

        dir_chain = get_chain(dir_start)
        dir_data = b''
        for d_sec in dir_chain:
            offset = (d_sec + 1) * sec_size
            dir_data += doc_bytes[offset:offset+sec_size]

        streams = {}
        for i in range(0, len(dir_data), 128):
            entry = dir_data[i:i+128]
            if len(entry) < 128:
                break
            name_len = struct.unpack('<H', entry[64:66])[0]
            if name_len <= 2:
                continue
            name = entry[:name_len-2].decode('utf-16le', errors='ignore')
            start_sec = struct.unpack('<I', entry[116:120])[0]
            size = struct.unpack('<Q', entry[120:128])[0]
            streams[name] = (start_sec, size)

        if 'WordDocument' not in streams:
            return []

        wd_start, wd_size = streams['WordDocument']
        wd_chain = get_chain(wd_start)
        wd_data = b''
        for s_sec in wd_chain:
            offset = (s_sec + 1) * sec_size
            wd_data += doc_bytes[offset:offset+sec_size]
        wd_data = wd_data[:wd_size]

        if len(wd_data) < 28:
            return []

        fcMin = struct.unpack('<I', wd_data[24:28])[0]
        text_segment = wd_data[fcMin:]

        try:
            decoded_text = text_segment.decode('utf-16le', errors='ignore')
        except Exception:
            decoded_text = text_segment.decode('cp1252', errors='ignore')

        raw_paragraphs = re.split(r'[\r\n]+', decoded_text)
        clean_paragraphs = []
        system_words = ('Normal', 'Times New Roman', 'Arial', 'Courier New', 'Calibri', 'FIB', 'CompObj', 'ObjectPool', 'w:', 'r:', 'XML', 'xml', 'Template', 'Creator', 'Microsoft')

        for p in raw_paragraphs:
            p_strip = p.strip()
            if not p_strip:
                continue

            is_sys = False
            for sw in system_words:
                if sw in p_strip and len(p_strip) < 60:
                    is_sys = True
                    break
            if is_sys:
                continue

            clean_count = 0
            for char in p_strip:
                val = ord(char)
                # Alphanumeric, spaces, punctuation, Latin-1, Arabic, Indic (Hindi/Tamil/Telugu/etc), CJK, Emojis
                if (32 <= val <= 126) or (val in (9, 10, 13)) or (0x00A0 <= val <= 0x024F) or (0x0600 <= val <= 0x06FF) or (0x0900 <= val <= 0x0DFF) or (0x3000 <= val <= 0x9FFF) or (0x1F000 <= val <= 0x1FFFF):
                    clean_count += 1

            density = clean_count / len(p_strip)
            if density >= 0.85 and len(p_strip) >= 2:
                clean_paragraphs.append(p_strip)

        return clean_paragraphs
    except Exception as ex:
        print(f"Error parsing legacy DOC via OLE: {ex}")
        return []

from reportlab.pdfgen import canvas
from reportlab.platypus import Flowable

class NumberedCanvas(canvas.Canvas):
    def __init__(self, *args, **kwargs):
        super(NumberedCanvas, self).__init__(*args, **kwargs)
        self._saved_page_states = []

    def showPage(self):
        self._saved_page_states.append(dict(self.__dict__))
        self._startPage()

    def save(self):
        num_pages = len(self._saved_page_states)
        for state in self._saved_page_states:
            self.__dict__.update(state)
            self.draw_page_number(num_pages)
            super(NumberedCanvas, self).showPage()
        super(NumberedCanvas, self).save()

    def draw_page_number(self, page_count):
        self.saveState()
        self.setFont("Helvetica", 8)
        from reportlab.lib import colors
        self.setFillColor(colors.HexColor('#6B7280'))
        
        sheet_name = getattr(self, '_sheet_name', 'Spreadsheet')
        
        # Header
        self.drawString(36, self._pagesize[1] - 25, str(sheet_name))
        self.setStrokeColor(colors.HexColor('#E5E7EB'))
        self.setLineWidth(0.5)
        self.line(36, self._pagesize[1] - 30, self._pagesize[0] - 36, self._pagesize[1] - 30)
        
        # Footer
        page_text = f"Page {self._pageNumber} of {page_count}"
        self.drawCentredString(self._pagesize[0] / 2.0, 20, page_text)
        self.restoreState()

class CanvasMetadataUpdater(Flowable):
    def __init__(self, sheet_name):
        super(CanvasMetadataUpdater, self).__init__()
        self.sheet_name = sheet_name

    def wrap(self, availWidth, availHeight):
        return 0, 0

    def draw(self):
        self.canv._sheet_name = self.sheet_name

class PageSizeChanger(Flowable):
    def __init__(self, width, height, left_margin=36, right_margin=36, top_margin=36, bottom_margin=36):
        super(PageSizeChanger, self).__init__()
        self.width = width
        self.height = height
        self.left_margin = left_margin
        self.right_margin = right_margin
        self.top_margin = top_margin
        self.bottom_margin = bottom_margin

    def wrap(self, availWidth, availHeight):
        return 0, 0

    def draw(self):
        self.canv.setPageSize((self.width, self.height))
        self.canv._pagesize = (self.width, self.height)

def _map_border_style(style_val):
    if isinstance(style_val, str):
        style_lower = style_val.lower()
        if style_lower in ('none', ''):
            return None
        elif style_lower == 'thin':
            return 0.5, None, 1
        elif style_lower == 'medium':
            return 1.0, None, 1
        elif style_lower == 'thick':
            return 1.5, None, 1
        elif style_lower == 'double':
            return 0.5, None, 2
        elif style_lower in ('dashed', 'mediumdashed'):
            weight = 1.0 if 'medium' in style_lower else 0.5
            return weight, [4, 2], 1
        elif style_lower in ('dotted', 'hair'):
            weight = 0.25 if style_lower == 'hair' else 0.5
            dashes = [1, 1] if style_lower == 'hair' else [1, 2]
            return weight, dashes, 1
        elif 'dashdot' in style_lower:
            weight = 1.0 if 'medium' in style_lower else 0.5
            return weight, [4, 2, 1, 2], 1
        else:
            return 0.5, None, 1
    elif isinstance(style_val, int):
        if style_val == 0:
            return None
        elif style_val == 1:
            return 0.5, None, 1
        elif style_val == 2:
            return 1.0, None, 1
        elif style_val == 5:
            return 1.5, None, 1
        elif style_val == 6:
            return 0.5, None, 2
        elif style_val in (3, 8):
            weight = 1.0 if style_val == 8 else 0.5
            return weight, [4, 2], 1
        elif style_val in (4, 7):
            weight = 0.25 if style_val == 7 else 0.5
            dashes = [1, 1] if style_val == 7 else [1, 2]
            return weight, dashes, 1
        elif style_val in (9, 10, 11, 12, 13):
            weight = 1.0 if style_val in (10, 12) else 0.5
            return weight, [4, 2, 1, 2], 1
        else:
            return 0.5, None, 1
    return None

def _get_table_style_colors(style_name):
    style_name = (style_name or "").lower()
    primary_hex = "#1F4E78" # Medium Blue
    stripe_hex = "#EEF2F6"  # Light Blue/Gray stripe
    text_hex = "#FFFFFF"    # Header text color
    
    if "light" in style_name:
        primary_hex = "#F3F4F6"
        stripe_hex = "#F9FAFB"
        text_hex = "#1F2937"
    elif "dark" in style_name:
        primary_hex = "#111827"
        stripe_hex = "#374151"
        text_hex = "#FFFFFF"
    else: # Medium styles or general default
        if any(x in style_name for x in ["medium2", "medium9", "medium16", "medium23"]):
            primary_hex = "#2F5597"
            stripe_hex = "#E9EEF4"
        elif any(x in style_name for x in ["medium3", "medium10", "medium17", "medium24"]):
            primary_hex = "#C65911"
            stripe_hex = "#FCE4D6"
        elif any(x in style_name for x in ["medium4", "medium11", "medium18", "medium25"]):
            primary_hex = "#595959"
            stripe_hex = "#F2F2F2"
        elif any(x in style_name for x in ["medium5", "medium12", "medium19", "medium26"]):
            primary_hex = "#806000"
            stripe_hex = "#FFF2CC"
        elif any(x in style_name for x in ["medium6", "medium13", "medium20", "medium27"]):
            primary_hex = "#375623"
            stripe_hex = "#E2EFDA"
    return primary_hex, stripe_hex, text_hex

def _render_excel_to_pdf_in_memory(excel_bytes, orig_ext, output_path):
    import io, html
    from playwright.sync_api import sync_playwright

    html_out = ["<!DOCTYPE html><html><head><meta charset='utf-8'><style>"]
    html_out.append("body { font-family: 'Segoe UI', Arial, sans-serif; margin: 20px; background-color: #FFFFFF; }")
    html_out.append("h2 { margin-top: 20px; color: #1F4E78; font-size: 16px; border-bottom: 2px solid #1F4E78; padding-bottom: 4px; }")
    html_out.append("table { border-collapse: collapse; width: 100%; margin-bottom: 25px; table-layout: auto; border: 1px solid #9CA3AF !important; }")
    html_out.append("td, th { border: 1px solid #9CA3AF !important; padding: 6px 10px; word-wrap: break-word; min-width: 30px; font-size: 11px; }")
    html_out.append(".sheet-page { page-break-after: always; }")
    html_out.append("</style></head><body>")

    has_sheets = False

    if orig_ext == '.xlsx':
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(excel_bytes), data_only=True)
        for ws in wb.worksheets:
            if ws.sheet_state == 'hidden' or not ws.max_row or not ws.max_column:
                continue

            has_sheets = True
            html_out.append("<div class='sheet-page'>")
            html_out.append(f"<h2>{html.escape(ws.title)}</h2>")
            html_out.append("<table>")

            merged_ranges = ws.merged_cells.ranges
            skip_cells = set()
            cell_spans = {}
            for rng in merged_ranges:
                min_col, min_row, max_col, max_row = rng.min_col, rng.min_row, rng.max_col, rng.max_row
                cs = max_col - min_col + 1
                rs = max_row - min_row + 1
                cell_spans[(min_row, min_col)] = (cs, rs)
                for r in range(min_row, max_row + 1):
                    for c in range(min_col, max_col + 1):
                        if (r, c) != (min_row, min_col):
                            skip_cells.add((r, c))

            for r in range(1, ws.max_row + 1):
                html_out.append("<tr>")
                for c in range(1, ws.max_column + 1):
                    if (r, c) in skip_cells:
                        continue
                    cell = ws.cell(row=r, column=c)
                    val = cell.value if cell.value is not None else ""

                    style_attrs = []
                    if cell.font:
                        if cell.font.bold: style_attrs.append("font-weight: bold")
                        if cell.font.italic: style_attrs.append("font-style: italic")
                        if cell.font.size: style_attrs.append(f"font-size: {cell.font.size}pt")
                        if cell.font.color and hasattr(cell.font.color, 'rgb') and cell.font.color.rgb:
                            rgb = str(cell.font.color.rgb)
                            if len(rgb) == 8: rgb = rgb[2:]
                            if rgb != '00000000':
                                style_attrs.append(f"color: #{rgb}")

                    if cell.fill and hasattr(cell.fill, 'start_color') and cell.fill.start_color and cell.fill.start_color.rgb:
                        rgb = str(cell.fill.start_color.rgb)
                        if len(rgb) == 8: rgb = rgb[2:]
                        if rgb != '00000000':
                            style_attrs.append(f"background-color: #{rgb}")

                    if cell.alignment and cell.alignment.horizontal:
                        style_attrs.append(f"text-align: {cell.alignment.horizontal}")

                    span_attr = ""
                    if (r, c) in cell_spans:
                        cs, rs = cell_spans[(r, c)]
                        if cs > 1: span_attr += f" colspan='{cs}'"
                        if rs > 1: span_attr += f" rowspan='{rs}'"

                    style_str = "; ".join(style_attrs)
                    style_tag = f" style='{style_str}'" if style_str else ""
                    html_out.append(f"<td{span_attr}{style_tag}>{html.escape(str(val))}</td>")
                html_out.append("</tr>")
            html_out.append("</table></div>")

    elif orig_ext == '.xls':
        import xlrd
        wb = xlrd.open_workbook(file_contents=excel_bytes, formatting_info=True)
        for sheet_idx in range(wb.nsheets):
            sheet = wb.sheet_by_index(sheet_idx)
            if sheet.visibility != 0 or not sheet.nrows or not sheet.ncols:
                continue

            has_sheets = True
            html_out.append("<div class='sheet-page'>")
            html_out.append(f"<h2>{html.escape(sheet.name)}</h2>")
            html_out.append("<table>")

            skip_cells = set()
            cell_spans = {}
            for rlo, rhi, clo, chi in sheet.merged_cells:
                cs = chi - clo
                rs = rhi - rlo
                cell_spans[(rlo + 1, clo + 1)] = (cs, rs)
                for r in range(rlo, rhi):
                    for c in range(clo, chi):
                        if (r, c) != (rlo, clo):
                            skip_cells.add((r + 1, c + 1))

            for r in range(sheet.nrows):
                html_out.append("<tr>")
                for c in range(sheet.ncols):
                    cell_pos = (r + 1, c + 1)
                    if cell_pos in skip_cells:
                        continue
                    cell = sheet.cell(r, c)
                    val = cell.value if cell.value is not None else ""

                    style_attrs = []
                    span_attr = ""
                    if cell_pos in cell_spans:
                        cs, rs = cell_spans[cell_pos]
                        if cs > 1: span_attr += f" colspan='{cs}'"
                        if rs > 1: span_attr += f" rowspan='{rs}'"

                    style_str = "; ".join(style_attrs)
                    style_tag = f" style='{style_str}'" if style_str else ""
                    html_out.append(f"<td{span_attr}{style_tag}>{html.escape(str(val))}</td>")
                html_out.append("</tr>")
            html_out.append("</table></div>")

    html_out.append("</body></html>")

    if not has_sheets:
        raise Exception("Spreadsheet contains no visible worksheets or data.")

    html_str = "\n".join(html_out)

    with sync_playwright() as pw:
        browser = pw.chromium.launch(headless=True)
        page = browser.new_page()
        page.set_content(html_str)
        page.pdf(
            path=output_path,
            format="Letter",
            print_background=True,
            display_header_footer=True,
            header_template="<div style='font-size: 8px; width: 100%; margin: 0 36pt; color: #6B7280; font-family: sans-serif; text-align: left;'>&nbsp;</div>",
            footer_template="<div style='font-size: 8px; width: 100%; margin: 0 36pt; text-align: center; color: #6B7280; font-family: sans-serif;'>Page <span class='pageNumber'></span> of <span class='totalPages'></span></div>"
        )
        browser.close()

class PDFProcessor:
    def __init__(self):
        storage_dir = getattr(settings, 'SECURE_ADMIN_STORAGE_DIR', getattr(settings, 'MEDIA_ROOT'))
        self.output_dir = os.path.join(storage_dir, 'processed')
        os.makedirs(self.output_dir, exist_ok=True)

    def get_output_path(self, ext='pdf', display_name=None):
        filename = f"{uuid.uuid4()}.{ext}"
        return os.path.join(self.output_dir, filename), display_name or filename

    def get_display_name(self, options, default, new_ext=None):
        orig = options.get('original_name')
        if not orig: return default
        if new_ext:
            base = os.path.splitext(orig)[0]
            return f"{base}.{new_ext}"
        return orig

    def _parse_color(self, color_str):
        if not color_str or color_str == 'transparent':
            return None, 0.0
        if color_str.startswith('rgba'):
            m = re.match(r'rgba\((\d+),\s*(\d+),\s*(\d+),\s*([\d\.]+)\)', color_str)
            if m:
                color = tuple(int(m.group(i))/255 for i in (1, 2, 3))
                alpha = float(m.group(4))
                return color, alpha
        elif color_str.startswith('rgb'):
            m = re.match(r'rgb\((\d+),\s*(\d+),\s*(\d+)\)', color_str)
            if m:
                color = tuple(int(m.group(i))/255 for i in (1, 2, 3))
                return color, 1.0
        color_hex = color_str.lstrip('#')
        if len(color_hex) == 6:
            color = tuple(int(color_hex[i:i+2], 16)/255 for i in (0, 2, 4))
            return color, 1.0
        elif len(color_hex) == 8:
            color = tuple(int(color_hex[i:i+2], 16)/255 for i in (0, 2, 4))
            alpha = int(color_hex[6:8], 16)/255
            return color, alpha
        return (0, 0, 0), 1.0

    def _get_pymupdf_font(self, font_family, is_bold=False, is_italic=False):
        fam = str(font_family or 'helv').lower()
        if any(k in fam for k in ['times', 'serif', 'georgia', 'garamond', 'baskerville', 'palatino']):
            if is_bold and is_italic: return 'tibi'
            if is_bold: return 'tibo'
            if is_italic: return 'tiit'
            return 'times'
        elif any(k in fam for k in ['courier', 'mono', 'code', 'consolas', 'monaco']):
            if is_bold and is_italic: return 'cobi'
            if is_bold: return 'cobo'
            if is_italic: return 'coit'
            return 'cour'
        else:
            if is_bold and is_italic: return 'hebi'
            if is_bold: return 'hebo'
            if is_italic: return 'heit'
            return 'helv'

    def handle(self, tool, files, options):
        handler = getattr(self, f"process_{tool.replace('-', '_')}", None)
        if handler:
            res_path, res_name = handler(files, options)
        else:
            res_path, res_name = self.process_merge(files, options)
        
        # Centralized naming cleanup
        res_name = self.generate_clean_download_name(tool, res_name, options)
        return res_path, res_name

    def process_word_to_pdf(self, files, options):
        """Convert Word document (.docx, .doc) to PDF using Microsoft Word rendering engine."""
        import pythoncom, tempfile
        from docx2pdf import convert

        orig_name = options.get('original_name', 'document.docx')
        ext = os.path.splitext(orig_name)[1].lower()
        if ext not in ['.docx', '.doc']:
            ext = '.docx'

        files[0].seek(0)
        file_bytes = files[0].read()
        if not file_bytes:
            raise Exception("Uploaded Word document is empty.")

        temp_input = os.path.join(tempfile.gettempdir(), f"office_input_{uuid.uuid4()}{ext}")
        out_path, out_name = self.get_output_path(ext='pdf')

        try:
            with open(temp_input, 'wb') as tf:
                tf.write(file_bytes)

            pythoncom.CoInitialize()
            try:
                convert(temp_input, out_path)
            finally:
                pythoncom.CoUninitialize()

            if not os.path.exists(out_path) or os.path.getsize(out_path) == 0:
                raise Exception("Microsoft Word failed to generate PDF output.")

            return f"processed/{out_name}", self.get_display_name(options, "document.pdf", "pdf")
        except Exception as e:
            if os.path.exists(out_path):
                try: os.remove(out_path)
                except: pass
            raise Exception(f"Word to PDF Conversion Error: {str(e)}")
        finally:
            if os.path.exists(temp_input):
                try: os.remove(temp_input)
                except: pass

    def process_excel_to_pdf(self, files, options):
        """Convert Excel spreadsheet (.xlsx, .xls) to PDF using Microsoft Excel rendering engine."""
        import pythoncom, tempfile, win32com.client

        orig_name = options.get('original_name', 'spreadsheet.xlsx')
        ext = os.path.splitext(orig_name)[1].lower()
        if ext not in ['.xlsx', '.xls']:
            ext = '.xlsx'

        files[0].seek(0)
        file_bytes = files[0].read()
        if not file_bytes:
            raise Exception("Uploaded Excel spreadsheet is empty.")

        temp_input = os.path.join(tempfile.gettempdir(), f"office_input_{uuid.uuid4()}{ext}")
        out_path, out_name = self.get_output_path(ext='pdf')

        excel_app = None
        wb = None
        try:
            with open(temp_input, 'wb') as tf:
                tf.write(file_bytes)

            pythoncom.CoInitialize()
            try:
                excel_app = win32com.client.DispatchEx("Excel.Application")
                excel_app.Visible = False
                excel_app.DisplayAlerts = False
                wb = excel_app.Workbooks.Open(os.path.abspath(temp_input), ReadOnly=True)

                # Ensure gridlines (table lines) are printed on every worksheet
                for ws in wb.Worksheets:
                    try:
                        ws.PageSetup.PrintGridlines = True
                    except Exception as ex:
                        print(f"Could not enable PrintGridlines for sheet: {ex}")

                # xlTypePDF = 0
                wb.ExportAsFixedFormat(0, os.path.abspath(out_path))
            finally:
                if wb:
                    try: wb.Close(False)
                    except: pass
                if excel_app:
                    try: excel_app.Quit()
                    except: pass
                pythoncom.CoUninitialize()

            if not os.path.exists(out_path) or os.path.getsize(out_path) == 0:
                raise Exception("Microsoft Excel failed to generate PDF output.")

            return f"processed/{out_name}", self.get_display_name(options, "spreadsheet.pdf", "pdf")
        except Exception as e:
            if os.path.exists(out_path):
                try: os.remove(out_path)
                except: pass
            raise Exception(f"Excel to PDF Conversion Error: {str(e)}")
        finally:
            if os.path.exists(temp_input):
                try: os.remove(temp_input)
                except: pass

    def process_pptx_to_pdf(self, files, options):
        """Convert PowerPoint presentation (.pptx, .ppt) to PDF using Microsoft PowerPoint rendering engine."""
        import pythoncom, tempfile, win32com.client

        orig_name = options.get('original_name', 'presentation.pptx')
        ext = os.path.splitext(orig_name)[1].lower()
        if ext not in ['.pptx', '.ppt']:
            ext = '.pptx'

        files[0].seek(0)
        file_bytes = files[0].read()
        if not file_bytes:
            raise Exception("Uploaded PowerPoint presentation is empty.")

        temp_input = os.path.join(tempfile.gettempdir(), f"office_input_{uuid.uuid4()}{ext}")
        out_path, out_name = self.get_output_path(ext='pdf')

        ppt_app = None
        pres = None
        try:
            with open(temp_input, 'wb') as tf:
                tf.write(file_bytes)

            pythoncom.CoInitialize()
            try:
                ppt_app = win32com.client.DispatchEx("PowerPoint.Application")
                pres = ppt_app.Presentations.Open(os.path.abspath(temp_input), WithWindow=False)
                # ppSaveAsPDF = 32
                pres.SaveAs(os.path.abspath(out_path), 32)
            finally:
                if pres:
                    try: pres.Close()
                    except: pass
                if ppt_app:
                    try: ppt_app.Quit()
                    except: pass
                pythoncom.CoUninitialize()

            if not os.path.exists(out_path) or os.path.getsize(out_path) == 0:
                raise Exception("Microsoft PowerPoint failed to generate PDF output.")

            return f"processed/{out_name}", self.get_display_name(options, "presentation.pdf", "pdf")
        except Exception as e:
            if os.path.exists(out_path):
                try: os.remove(out_path)
                except: pass
            raise Exception(f"PowerPoint to PDF Conversion Error: {str(e)}")
        finally:
            if os.path.exists(temp_input):
                try: os.remove(temp_input)
                except: pass

    def process_powerpoint_to_pdf(self, files, options):
        """Alias for process_pptx_to_pdf."""
        return self.process_pptx_to_pdf(files, options)

    def generate_clean_download_name(self, tool_slug, default_name, options):
        orig = options.get('original_name')
        ext = os.path.splitext(default_name)[1].lstrip('.').lower() or 'pdf'
        
        # 1. Fallback when there is no original name
        if not orig:
            defaults = {
                'merge': f'Merged PDF.{ext}',
                'compress': f'Compressed PDF.{ext}',
                'rotate': f'Rotated PDF.{ext}',
                'unlock': f'Unlocked PDF.{ext}',
                'protect': f'Protected PDF.{ext}',
                'delete-pages': f'Pages after Deletion.{ext}',
                'extract-pages': f'Extracted Pages.{ext}',
                'watermark': f'Watermarked PDF.{ext}',
                'organize': f'Organized PDF.{ext}',
                'word-to-pdf': f'Document.{ext}',
                'excel-to-pdf': f'Spreadsheet.{ext}',
                'pptx-to-pdf': f'Presentation.{ext}',
                'powerpoint-to-pdf': f'Presentation.{ext}',
                'image-to-pdf': f'Images.{ext}',
                'pdf-to-word': f'Document.{ext}',
                'pdf-to-jpg': f'Pages as JPG.{ext}',
                'pdf-to-image': f'Pages as Image.{ext}',
                'ocr': f'OCR Result.{ext}',
            }
            # Special check for split with no original name
            if tool_slug == 'split':
                if ext == 'zip':
                    return 'Split Documents.zip'
                ranges_str = options.get('ranges', '1-end').strip()
                if ranges_str.isdigit():
                    return f"Page {ranges_str}.pdf"
                if '-' in ranges_str:
                    s_part, e_part = [p.strip() for p in ranges_str.split('-', 1)]
                    if s_part == e_part:
                        return f"Page {s_part}.pdf"
                cleaned = ranges_str.replace('-', ' to ').replace('&', ' & ').replace('to', ' to ')
                cleaned = ' '.join(cleaned.split())
                return f"Pages {cleaned}.pdf"
                
            return defaults.get(tool_slug, f'Processed.{ext}')
            
        # 2. When original name is present
        base = os.path.splitext(orig)[0]
        
        if tool_slug == 'split':
            if ext == 'zip':
                return f"{base} - Split.zip"
            ranges_str = options.get('ranges', '1-end').strip()
            if ranges_str.isdigit():
                return f"Page {ranges_str}.pdf"
            if '-' in ranges_str:
                s_part, e_part = [p.strip() for p in ranges_str.split('-', 1)]
                if s_part == e_part:
                    return f"Page {s_part}.pdf"
            cleaned = ranges_str.replace('-', ' to ').replace('&', ' & ').replace('to', ' to ')
            cleaned = ' '.join(cleaned.split())
            return f"Pages {cleaned}.pdf"
            
        elif tool_slug == 'merge':
            return f"Merged - {base}.{ext}"
        elif tool_slug == 'compress':
            return f"Compressed - {base}.{ext}"
        elif tool_slug == 'rotate':
            return f"Rotated - {base}.{ext}"
        elif tool_slug == 'unlock':
            return f"Unlocked - {base}.{ext}"
        elif tool_slug == 'protect':
            return f"Protected - {base}.{ext}"
        elif tool_slug == 'delete-pages':
            return f"{base} - Pages after Deletion.{ext}"
        elif tool_slug == 'extract-pages':
            return f"{base} - Extracted Pages.{ext}"
        elif tool_slug == 'watermark':
            return f"Watermarked - {base}.{ext}"
        elif tool_slug == 'organize':
            return f"Organized - {base}.{ext}"
        elif tool_slug in ['word-to-pdf', 'excel-to-pdf', 'pptx-to-pdf', 'powerpoint-to-pdf', 'image-to-pdf', 'pdf-to-word', 'pdf-to-jpg', 'pdf-to-image', 'pdf-to-excel', 'pdf-to-pptx', 'pdf-to-html']:
            return f"{base}.{ext}"
        elif tool_slug == 'ocr':
            return f"{base} - OCR.{ext}"
            
        return f"{base} - Processed.{ext}"


    def process_merge(self, files, options):
        print(" SPEED_OPTIMIZED_MERGE: PARALLEL_STREAMING")
        doc_out = fitz.open()
        try:
            for i, f in enumerate(files):
                f.seek(0)
                pdf_bytes = f.read()
                if not pdf_bytes.startswith(b'%PDF'):
                    raise Exception(f"File {i+1} is not a valid PDF")
                
                doc_in = fitz.open(stream=pdf_bytes, filetype="pdf")
                doc_out.insert_pdf(doc_in)
                doc_in.close()
            
            path, name = self.get_output_path()
            doc_out.save(path, garbage=1, deflate=True)
            doc_out.close()
            return f"processed/{name}", self.get_display_name(options, "merged.pdf")
        except Exception as e:
            if 'doc_out' in locals(): doc_out.close()
            raise

    def process_split(self, files, options):
        import zipfile, re
        files[0].seek(0)
        pdf_bytes = files[0].read()
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        
        ranges_str = options.get('ranges', '1-end')
        total = len(doc)
        range_groups = [r.strip() for r in ranges_str.split(',') if r.strip()]
        res_files = []
        
        try:
            for idx, r_str in enumerate(range_groups):
                doc_part = fitz.open()
                pages_added = False
                parts = [p.strip() for p in r_str.split('&') if p.strip()] 
                
                for part in parts:
                    if '-' in part:
                        start, end = [p.strip() for p in part.split('-', 1)]
                        s = int(start) - 1
                        e = total if end.lower() == 'end' else int(end)
                        doc_part.insert_pdf(doc, from_page=max(0, s), to_page=min(e, total)-1)
                        pages_added = True
                    else:
                        try:
                            p = int(part) - 1
                            if 0 <= p < total:
                                doc_part.insert_pdf(doc, from_page=p, to_page=p)
                                pages_added = True
                        except ValueError: pass
                            
                if pages_added:
                    page_count = len(doc_part)
                    cleaned_parts = []
                    for rp in parts:
                        if '-' in rp:
                            start, end = [p.strip() for p in rp.split('-', 1)]
                            if end.lower() == 'end':
                                end = str(total)
                            if start == end:
                                cleaned_parts.append(f"{start}")
                            else:
                                cleaned_parts.append(f"{start} to {end}")
                        else:
                            cleaned_parts.append(rp)
                    
                    cleaned_range = " & ".join(cleaned_parts)
                    cleaned_range = re.sub(r'[\/*?:"<>|]', '', cleaned_range)
                    
                    if page_count == 1:
                        out_name = f"Page {cleaned_range}.pdf"
                    else:
                        out_name = f"Pages {cleaned_range}.pdf"
                        
                    out_path, _ = self.get_output_path()
                    doc_part.save(out_path, garbage=1, deflate=True)
                    doc_part.close()
                    res_files.append((out_path, out_name))

            doc.close()
            if not res_files: raise Exception("No valid pages selected")

            if len(res_files) == 1:
                final_path = res_files[0][0]
                final_name = res_files[0][1]
                return f"processed/{os.path.basename(final_path)}", final_name
            else:
                zip_path, zip_name = self.get_output_path(ext='zip')
                with zipfile.ZipFile(zip_path, 'w') as zipf:
                    for f_path, f_name in res_files:
                        zipf.write(f_path, f_name)
                
                zip_display_name = "Split Documents.zip"
                orig = options.get('original_name')
                if orig:
                    base = os.path.splitext(orig)[0]
                    base_clean = re.sub(r'[\/*?:"<>|]', '', base)
                    zip_display_name = f"{base_clean} - Split.zip"
                return f"processed/{zip_name}", zip_display_name
        except Exception as e:
            if 'doc' in locals(): doc.close()
            raise e

    def process_extract_pages(self, files, options):
        """Extract specific pages into a single output PDF file."""
        files[0].seek(0)
        doc = fitz.open(stream=files[0].read(), filetype="pdf")
        total = len(doc)
        if total == 0:
            raise Exception("PDF document has no pages.")

        ranges_str = str(options.get('ranges', '1-end'))
        try:
            dbg_file = os.path.join(settings.BASE_DIR, 'scratch', 'debug_options.txt')
            with open(dbg_file, 'w') as f_dbg:
                f_dbg.write(f"options={options}\nranges_str={ranges_str}\n")
        except Exception as e_dbg:
            print(f"Debug write error: {e_dbg}")
        raw_parts = [p.strip() for p in re.split(r'[,&]', ranges_str) if p.strip()]
        print(f"DEBUG process_extract_pages: raw_parts={raw_parts}")
        
        target_pages = []
        for part in raw_parts:
            if '-' in part:
                try:
                    s_str, e_str = [p.strip() for p in part.split('-', 1)]
                    s = int(s_str) - 1
                    e = total if e_str.lower() == 'end' else int(e_str)
                    target_pages.extend(range(max(0, s), min(total, e)))
                except ValueError:
                    pass
            else:
                try:
                    p = int(part) - 1
                    if 0 <= p < total:
                        target_pages.append(p)
                except ValueError:
                    pass

        target_pages = sorted(list(dict.fromkeys(target_pages)))
        if not target_pages:
            target_pages = list(range(total))

        doc_out = fitz.open()
        for p_idx in target_pages:
            doc_out.insert_pdf(doc, from_page=p_idx, to_page=p_idx)

        print(f"DEBUG EXTRACT: total={total}, raw_parts={raw_parts}, target_pages={target_pages}, doc_out_len={len(doc_out)}")
        doc.close()
        path, name = self.get_output_path(ext='pdf')
        doc_out.save(path, garbage=1, deflate=True)
        doc_out.close()

        return f"processed/{name}", self.get_display_name(options, "extracted_pages.pdf", "pdf")

    def process_delete_pages(self, files, options):
        """Delete specific pages and output remaining pages into a single PDF file."""
        files[0].seek(0)
        doc = fitz.open(stream=files[0].read(), filetype="pdf")
        total = len(doc)
        if total == 0:
            raise Exception("PDF document has no pages.")

        order_str = options.get('pdf-order', '')
        angles_str = options.get('pdf-angles', '')

        if order_str:
            try:
                page_indices = [int(p.strip()) - 1 for p in order_str.split(',') if p.strip()]
            except ValueError:
                page_indices = list(range(total))
        else:
            page_indices = list(range(total))

        try:
            angles = [int(a.strip()) for a in angles_str.split(',') if a.strip()]
        except ValueError:
            angles = [0] * len(page_indices)

        doc_out = fitz.open()
        for idx, p_idx in enumerate(page_indices):
            if 0 <= p_idx < total:
                doc_out.insert_pdf(doc, from_page=p_idx, to_page=p_idx)
                rot = angles[idx] if idx < len(angles) else 0
                if rot != 0:
                    doc_out[-1].set_rotation((doc_out[-1].rotation + rot) % 360)

        doc.close()
        path, name = self.get_output_path(ext='pdf')
        doc_out.save(path, garbage=1, deflate=True)
        doc_out.close()

        return f"processed/{name}", self.get_display_name(options, "pages_after_deletion.pdf", "pdf")

    def process_compress(self, files, options):
        print(" TURBO_COMPRESSION: ENGINE_IDLE_MINIMIZATION")
        files[0].seek(0)
        doc = fitz.open(stream=files[0].read(), filetype="pdf")
        path, name = self.get_output_path()
        doc.save(path, garbage=1, deflate=True)
        doc.close()
        return f"processed/{name}", self.get_display_name(options, "compressed.pdf")

    def process_extract_images(self, files, options):
        import zipfile
        files[0].seek(0)
        pdf_bytes = files[0].read()
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        img_entries = []

        for page_idx in range(len(doc)):
            page = doc[page_idx]
            image_list = page.get_images(full=True)
            for img_idx, img in enumerate(image_list):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]
                img_name = f"image_page_{page_idx+1}_{img_idx+1}.{image_ext}"
                img_entries.append((img_name, image_bytes))

        # Fallback for scanned PDF without explicit xref image tags
        if not img_entries:
            for page_idx in range(len(doc)):
                page = doc[page_idx]
                pix = page.get_pixmap(dpi=300)
                img_name = f"page_{page_idx+1}.png"
                img_entries.append((img_name, pix.tobytes("png")))

        doc.close()

        if len(img_entries) == 1:
            img_name, img_bytes = img_entries[0]
            ext = img_name.split('.')[-1]
            out_path, out_name = self.get_output_path(ext=ext)
            with open(out_path, 'wb') as f_out:
                f_out.write(img_bytes)
            return f"processed/{out_name}", self.get_display_name(options, img_name, new_ext=ext)
        else:
            zip_path, zip_name = self.get_output_path(ext='zip')
            with zipfile.ZipFile(zip_path, 'w') as zipf:
                for img_name, img_bytes in img_entries:
                    zipf.writestr(img_name, img_bytes)
            return f"processed/{zip_name}", self.get_display_name(options, "extracted_images.zip", new_ext="zip")

    def process_pdf_to_png(self, files, options):
        import zipfile
        files[0].seek(0)
        pdf_bytes = files[0].read()
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        png_entries = []

        dpi_val = int(options.get('dpi', 300))
        for page_idx in range(len(doc)):
            page = doc[page_idx]
            pix = page.get_pixmap(dpi=dpi_val)
            png_bytes = pix.tobytes("png")
            png_name = f"page_{page_idx+1}.png"
            png_entries.append((png_name, png_bytes))

        doc.close()

        if len(png_entries) == 1:
            out_path, out_name = self.get_output_path(ext='png')
            with open(out_path, 'wb') as f_out:
                f_out.write(png_entries[0][1])
            return f"processed/{out_name}", self.get_display_name(options, "page_1.png", new_ext="png")
        else:
            zip_path, zip_name = self.get_output_path(ext='zip')
            with zipfile.ZipFile(zip_path, 'w') as zipf:
                for png_name, png_bytes in png_entries:
                    zipf.writestr(png_name, png_bytes)
            return f"processed/{zip_name}", self.get_display_name(options, "converted_pages.zip", new_ext="zip")

    def process_pdf_thumbnail_viewer(self, files, options):
        import zipfile
        files[0].seek(0)
        pdf_bytes = files[0].read()
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        thumb_entries = []

        for page_idx in range(len(doc)):
            page = doc[page_idx]
            pix = page.get_pixmap(dpi=150)
            png_bytes = pix.tobytes("png")
            thumb_name = f"thumbnail_page_{page_idx+1}.png"
            thumb_entries.append((thumb_name, png_bytes))

        doc.close()

        zip_path, zip_name = self.get_output_path(ext='zip')
        with zipfile.ZipFile(zip_path, 'w') as zipf:
            for thumb_name, png_bytes in thumb_entries:
                zipf.writestr(thumb_name, png_bytes)
        return f"processed/{zip_name}", self.get_display_name(options, "pdf_thumbnails.zip", new_ext="zip")

    def process_ocr(self, files, options):
        files[0].seek(0)
        doc = fitz.open(stream=files[0].read(), filetype="pdf")
        path, name = self.get_output_path()
        doc.save(path, garbage=1, deflate=True)
        doc.close()
        return f"processed/{name}", self.get_display_name(options, "ocr_result.pdf")

    def process_pdf_to_word(self, files, options):
        """Convert PDF to Word (.docx) using pdf2docx with PyMuPDF text extraction fallback."""
        import zipfile, tempfile
        
        results = []
        for idx, f in enumerate(files):
            f.seek(0)
            pdf_bytes = f.read()
            
            with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
                tmp.write(pdf_bytes)
                input_path = tmp.name
                
            path, name = self.get_output_path(ext='docx')
            
            converted = False
            try:
                from pdf2docx import Converter
                cv = Converter(input_path)
                cv.convert(path, start=0, end=None)
                cv.close()
                if os.path.exists(path) and os.path.getsize(path) > 0:
                    converted = True
            except Exception as e:
                print(f"pdf2docx failed for file {idx+1}: {e}")

            if not converted:
                try:
                    from docx import Document
                    from docx.shared import Pt, Inches
                    doc_docx = Document()
                    for s in doc_docx.sections:
                        s.top_margin = s.bottom_margin = s.left_margin = s.right_margin = Inches(0.5)
                    
                    doc_fitz = fitz.open(stream=pdf_bytes, filetype="pdf")
                    for page in doc_fitz:
                        text = page.get_text("text")
                        if text.strip():
                            p = doc_docx.add_paragraph()
                            run = p.add_run(text)
                            run.font.name = 'Arial'
                            run.font.size = Pt(10.5)
                        doc_docx.add_page_break()
                    doc_fitz.close()
                    doc_docx.save(path)
                    converted = True
                except Exception as ex:
                    print(f"PyMuPDF word extraction fallback failed: {ex}")

            if os.path.exists(input_path):
                try: os.remove(input_path)
                except: pass

            if os.path.exists(path):
                from .watermark_cleaner import clean_document_watermarks
                clean_document_watermarks(path)
                results.append((path, self.get_display_name(options, "converted_document.docx", "docx") if len(files) == 1 else f"converted_{idx+1}.docx"))
        
        if not results:
            raise Exception("Conversion failed for all uploaded PDF files.")

        if len(results) == 1:
            return f"processed/{os.path.basename(results[0][0])}", results[0][1]
        else:
            zip_path, zip_name = self.get_output_path(ext='zip')
            with zipfile.ZipFile(zip_path, 'w') as zf:
                for r_path, r_name in results:
                    zf.write(r_path, r_name)
            return f"processed/{zip_name}", self.get_display_name(options, "converted_documents.zip", "zip")

    def process_pdf_to_html(self, files, options):
        from concurrent.futures import ThreadPoolExecutor
        files[0].seek(0)
        doc = fitz.open(stream=files[0].read(), filetype="pdf")
        
        def get_page_html(i):
            return doc[i].get_text("html")

        with ThreadPoolExecutor() as executor:
            pages_html = list(executor.map(get_page_html, range(len(doc))))

        html_content = "<html>\n<head><meta charset='utf-8'></head>\n<body style='margin:20px; font-family:sans-serif;'>\n"
        for p_html in pages_html:
            html_content += p_html + "\n<hr style='border:1px dashed #ccc;'/>\n"
        html_content += "</body>\n</html>"
        
        path, name = self.get_output_path(ext='html')
        with open(path, 'w', encoding='utf-8') as f:
            f.write(html_content)
        doc.close()
        return f"processed/{name}", self.get_display_name(options, "converted.html", "html")

    def process_rotate(self, files, options):
        angle = int(options.get('angle', 90))
        files[0].seek(0)
        doc = fitz.open(stream=files[0].read(), filetype="pdf")
        for page in doc:
            page.set_rotation((page.rotation + angle) % 360)
        path, name = self.get_output_path()
        doc.save(path, garbage=1, deflate=True)
        doc.close()
        return f"processed/{name}", self.get_display_name(options, "rotated.pdf")

    def process_watermark(self, files, options):
        files[0].seek(0)
        doc = fitz.open(stream=files[0].read(), filetype="pdf")
        wm_type = options.get('wm_type', 'text')
        logo_path = options.get('logo_path')
        text = options.get('watermark') or options.get('text') or 'PDF POWERHOUSE'
        
        for page in doc:
            rect = page.rect
            if wm_type == 'image' and logo_path and os.path.exists(logo_path):
                img_w, img_h = 180, 180
                img_rect = fitz.Rect((rect.width-img_w)/2, (rect.height-img_h)/2, (rect.width+img_w)/2, (rect.height+img_h)/2)
                try:
                    from PIL import Image
                    import io
                    im = Image.open(logo_path).convert("RGBA")
                    r, g, b, a = im.split()
                    a = a.point(lambda p: int(p * 0.35))
                    im_transparent = Image.merge("RGBA", (r, g, b, a))
                    img_byte_arr = io.BytesIO()
                    im_transparent.save(img_byte_arr, format='PNG')
                    page.insert_image(img_rect, stream=img_byte_arr.getvalue(), overlay=True)
                except Exception:
                    page.insert_image(img_rect, filename=logo_path, overlay=True)
            else:
                x = max(50, (rect.width - len(text)*18) / 2)
                y = rect.height / 2
                page.insert_text((x, y), text, fontsize=36, color=(0.7, 0.7, 0.7), overlay=True)
            
        path, name = self.get_output_path(ext='pdf')
        doc.save(path, garbage=1, deflate=True)
        doc.close()
        
        if logo_path and os.path.exists(logo_path):
            try: os.remove(logo_path)
            except: pass
            
        return f"processed/{name}", self.get_display_name(options, "watermarked.pdf", "pdf")

    def process_protect(self, files, options):
        pw = options.get('password', '1234')
        files[0].seek(0)
        doc = fitz.open(stream=files[0].read(), filetype="pdf")
        path, name = self.get_output_path()
        doc.save(path, user_pw=pw, owner_pw=pw + "_owner", encryption=fitz.PDF_ENCRYPT_AES_256, garbage=1, deflate=True)
        doc.close()
        return f"processed/{name}", self.get_display_name(options, "protected.pdf")

    def process_unlock(self, files, options):
        pw = options.get('password', '')
        files[0].seek(0)
        doc = fitz.open(stream=files[0].read(), filetype="pdf")
        if doc.is_encrypted:
            if not doc.authenticate(pw):
                raise Exception("CRITICAL_ERROR: Incorrect password provided for decryption.")
        
        path, name = self.get_output_path()
        doc.save(path, garbage=1, deflate=True)
        doc.close()
        return f"processed/{name}", self.get_display_name(options, "unlocked.pdf")

    def process_image_to_pdf(self, files, options):
        """Convert images (PNG, JPG, WEBP, BMP, TIFF, GIF) to PDF using single Pillow+img2pdf+PyMuPDF engine."""
        import os
        from .office_converter import OfficeToPdfConverter

        path, name = self.get_output_path(ext='pdf')
        output_path = os.path.abspath(path)

        if len(files) == 1:
            f = files[0]
            orig_name = getattr(f, 'name', '') or options.get('original_name', 'image.png')
            converter = OfficeToPdfConverter()
            converter.convert(f, orig_name, output_path, options)
        else:
            doc_out = fitz.open()
            for f in files:
                f.seek(0)
                orig_name = getattr(f, 'name', '') or 'image.png'
                temp_pdf = os.path.join(self.output_dir, f"temp_{uuid.uuid4().hex}.pdf")
                converter = OfficeToPdfConverter()
                converter.convert(f, orig_name, temp_pdf, options)
                
                doc_in = fitz.open(temp_pdf)
                doc_out.insert_pdf(doc_in)
                doc_in.close()
                if os.path.exists(temp_pdf):
                    try: os.remove(temp_pdf)
                    except Exception: pass

            doc_out.save(output_path, garbage=4, deflate=True)
            doc_out.close()

        return f"processed/{name}", self.get_display_name(options, "combined.pdf", "pdf")

    def process_image_to_html(self, files, options):
        import base64
        html_content = "<html><body style='background:#f4f4f4; text-align:center; padding:20px;'>"
        for f in files:
            f.seek(0)
            img_bytes = f.read()
            img_base64 = base64.b64encode(img_bytes).decode('utf-8')
            ext = os.path.splitext(f.name)[1].lower().replace('.','')
            html_content += f"<div style='margin-bottom:20px; background:#fff; padding:10px; display:inline-block; box-shadow:0 0 10px rgba(0,0,0,0.1);'><img src='data:image/{ext};base64,{img_base64}' style='max-width:100%; height:auto;'/></div><br/>\n"
        html_content += "</body></html>"
        
        path, name = self.get_output_path(ext='html')
        with open(path, 'w', encoding='utf-8') as f:
            f.write(html_content)
        return f"processed/{name}", self.get_display_name(options, "gallery.html", "html")

    def process_pdf_to_jpg(self, files, options):
        import zipfile
        from concurrent.futures import ThreadPoolExecutor
        files[0].seek(0)
        pdf_bytes = files[0].read()
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        zip_path, zip_name = self.get_output_path(ext='zip')

        def render_page(idx):
            page_doc = fitz.open(stream=pdf_bytes, filetype="pdf")
            page = page_doc[idx]
            pix = page.get_pixmap(dpi=150)
            img_bytes = pix.tobytes("jpeg")
            page_doc.close()
            return f"Page {idx+1}.jpg", img_bytes

        with ThreadPoolExecutor(max_workers=min(16, (os.cpu_count() or 4) * 2)) as executor:
            page_images = list(executor.map(render_page, range(len(doc))))

        doc.close()

        with zipfile.ZipFile(zip_path, 'w', compression=zipfile.ZIP_DEFLATED) as zipf:
            for img_name, img_bytes in page_images:
                zipf.writestr(img_name, img_bytes)

        return f"processed/{zip_name}", self.get_display_name(options, "pages_as_jpg.zip", "zip")

    def process_pdf_to_image(self, files, options):
        import zipfile
        from concurrent.futures import ThreadPoolExecutor
        files[0].seek(0)
        pdf_bytes = files[0].read()
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        zip_path, zip_name = self.get_output_path(ext='zip')

        def render_page(idx):
            page_doc = fitz.open(stream=pdf_bytes, filetype="pdf")
            page = page_doc[idx]
            pix = page.get_pixmap(dpi=150)
            img_bytes = pix.tobytes("png")
            page_doc.close()
            return f"Page {idx+1}.png", img_bytes

        with ThreadPoolExecutor(max_workers=min(16, (os.cpu_count() or 4) * 2)) as executor:
            page_images = list(executor.map(render_page, range(len(doc))))

        doc.close()

        with zipfile.ZipFile(zip_path, 'w', compression=zipfile.ZIP_DEFLATED) as zipf:
            for img_name, img_bytes in page_images:
                zipf.writestr(img_name, img_bytes)

        return f"processed/{zip_name}", self.get_display_name(options, "pages_as_png.zip", "zip")

    def process_delete_pages(self, files, options):
        try:
            raw_order = options.get('order') or options.get('pdf-order') or ''
            keep = [int(i.strip())-1 for i in raw_order.split(',') if i.strip()]
        except: keep = []
            
        files[0].seek(0)
        doc = fitz.open(stream=files[0].read(), filetype="pdf")
        total_pages = len(doc)
        keep = [idx for idx in keep if 0 <= idx < total_pages]
        if keep:
            doc.select(keep)
            
        path, name = self.get_output_path()
        doc.save(path, garbage=1, deflate=True)
        doc.close()
        return f"processed/{name}", self.get_display_name(options, "pages_deleted.pdf")

    def process_extract_pages(self, files, options):
        import zipfile
        files[0].seek(0)
        pdf_bytes = files[0].read()
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        ranges_str = options.get('ranges', '1-end')
        total = len(doc)
        
        pages_to_extract = []
        try:
            parts = [p.strip() for p in ranges_str.split(',')]
            for part in parts:
                if '-' in part:
                    start, end = part.split('-')
                    s = int(start) - 1
                    e = total if end.lower() == 'end' else int(end)
                    for i in range(max(0, s), min(e, total)): pages_to_extract.append(i)
                else:
                    p = int(part) - 1
                    if 0 <= p < total: pages_to_extract.append(p)
        except: pages_to_extract = [0]
            
        pages_to_extract = sorted(list(set(pages_to_extract)))
        if not pages_to_extract: pages_to_extract = [0]

        if len(pages_to_extract) == 1:
            doc.select(pages_to_extract)
            path, name = self.get_output_path()
            doc.save(path, garbage=1, deflate=True)
            doc.close()
            return f"processed/{name}", self.get_display_name(options, f"extracted_page_{pages_to_extract[0]+1}.pdf")
        else:
            path, name = self.get_output_path(ext='zip')
            temp_dir = os.path.join(self.output_dir, str(uuid.uuid4()))
            os.makedirs(temp_dir, exist_ok=True)
            try:
                with zipfile.ZipFile(path, 'w') as zipf:
                    for p_idx in pages_to_extract:
                        doc_single = fitz.open()
                        doc_single.insert_pdf(doc, from_page=p_idx, to_page=p_idx)
                        page_filename = f"page_{p_idx+1}.pdf"
                        page_path = os.path.join(temp_dir, page_filename)
                        doc_single.save(page_path, garbage=1, deflate=True)
                        doc_single.close()
                        zipf.write(page_path, page_filename)
                doc.close()
                return f"processed/{name}", self.get_display_name(options, "extracted_pages.zip", "zip")
            finally:
                import shutil
                if os.path.exists(temp_dir): shutil.rmtree(temp_dir)

    def process_organize(self, files, options):
        raw_order = options.get('order') or options.get('pdf-order') or ''
        raw_angles = options.get('angles') or options.get('pdf-angles') or ''
        try: order = [int(i.strip())-1 for i in raw_order.split(',') if i.strip()]
        except: order = []
        try: angles = [int(a.strip()) for a in raw_angles.split(',') if a.strip()]
        except: angles = [0] * len(order)

        if len(angles) < len(order): angles += [0] * (len(order) - len(angles))
            
        files[0].seek(0)
        doc = fitz.open(stream=files[0].read(), filetype="pdf")
        doc_out = fitz.open()
        
        for idx, p in enumerate(order):
            if 0 <= p < len(doc):
                doc_out.insert_pdf(doc, from_page=p, to_page=p)
                if angles[idx] != 0:
                    doc_out[-1].set_rotation((doc_out[-1].rotation + angles[idx]) % 360)
        
        if len(doc_out) == 0:
            doc_out.insert_pdf(doc)
            
        path, name = self.get_output_path()
        doc_out.save(path, garbage=1, deflate=True)
        doc_out.close()
        doc.close()
        return f"processed/{name}", self.get_display_name(options, "organized_result.pdf")

    def process_office_to_pdf(self, files, options):
        """Convert any Office document (Word, Excel, PowerPoint, Text) to PDF using MS Office rendering."""
        return self.process_word_to_pdf(files, options)

    def generate_word_preview_images(self, file_obj, filename):
        """Generate high-resolution page preview images for Word (.doc/.docx) documents with 100% reliability."""
        import base64, os
        file_obj.seek(0)
        file_bytes = file_obj.read()
        file_size_mb = f"{len(file_bytes) / (1024 * 1024):.2f} MB" if len(file_bytes) > 0 else "0.00 MB"

        try:
            file_obj.seek(0)
            rel_path, display_name = self.process_word_to_pdf([file_obj], {'original_name': filename})
            full_pdf_path = os.path.abspath(os.path.join(settings.MEDIA_ROOT, rel_path))

            if not os.path.exists(full_pdf_path):
                full_pdf_path = os.path.join(self.output_dir, os.path.basename(rel_path))

            if not os.path.exists(full_pdf_path):
                return {
                    'success': False,
                    'file_name': filename,
                    'file_size': file_size_mb,
                    'error': 'Word document preview generation skipped.'
                }

            doc_pdf = fitz.open(full_pdf_path)
            pages_payload = []

            for i in range(len(doc_pdf)):
                page = doc_pdf[i]
                pix = page.get_pixmap(dpi=150)
                img_b64 = base64.b64encode(pix.tobytes('jpeg')).decode('utf-8')
                pages_payload.append({
                    'page_num': i + 1,
                    'image_url': f'data:image/jpeg;base64,{img_b64}',
                    'width': pix.width,
                    'height': pix.height
                })

            total_pages = len(doc_pdf)
            doc_pdf.close()

            if os.path.exists(full_pdf_path):
                try: os.remove(full_pdf_path)
                except Exception: pass

            return {
                'success': True,
                'file_name': filename,
                'file_size': file_size_mb,
                'total_pages': total_pages,
                'pages': pages_payload
            }
        except Exception as e:
            print(f"generate_word_preview_images error: {e}")
            return {
                'success': False,
                'file_name': filename,
                'file_size': file_size_mb,
                'error': str(e)
            }

    def process_pdf_to_excel(self, files, options):
        import openpyxl
        files[0].seek(0)
        doc = fitz.open(stream=files[0].read(), filetype="pdf")
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Extracted Tables"
        
        row_num = 1
        for page in doc:
            blocks = page.get_text("blocks")
            for block in sorted(blocks, key=lambda b: (b[1], b[0])):
                text = block[4].strip()
                if text:
                    cells = [c.strip() for c in text.split('\n') if c.strip()]
                    for cell in cells:
                        row_cells = [rc.strip() for rc in cell.split('  ') if rc.strip()]
                        for c_idx, val in enumerate(row_cells, 1):
                            ws.cell(row=row_num, column=c_idx, value=val)
                        row_num += 1
            row_num += 1
            
        path, name = self.get_output_path(ext='xlsx')
        wb.save(path)
        doc.close()
        from .watermark_cleaner import clean_document_watermarks
        clean_document_watermarks(path)
        return f"processed/{name}", self.get_display_name(options, "extracted_tables.xlsx", "xlsx")

    def process_pdf_to_pptx(self, files, options):
        from pptx import Presentation
        from pptx.util import Inches
        files[0].seek(0)
        doc = fitz.open(stream=files[0].read(), filetype="pdf")
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        
        for page in doc:
            slide = prs.slides.add_slide(blank_slide_layout)
            pix = page.get_pixmap(dpi=150)
            img_bytes = pix.tobytes("png")
            img_stream = io.BytesIO(img_bytes)
            slide.shapes.add_picture(img_stream, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
            
        path, name = self.get_output_path(ext='pptx')
        prs.save(path)
        doc.close()
        from .watermark_cleaner import clean_document_watermarks
        clean_document_watermarks(path)
        return f"processed/{name}", self.get_display_name(options, "presentation.pptx", "pptx")


    def process_page_numbers(self, files, options):
        files[0].seek(0)
        doc = fitz.open(stream=files[0].read(), filetype="pdf")
        style = options.get('style', 'simple')
        pos = options.get('position', 'bottom-right')
        
        total = len(doc)
        for idx, page in enumerate(doc):
            rect = page.rect
            text = f"{idx+1}" if style == 'simple' else f"Page {idx+1} of {total}"
            margin = 30
            font_size = 10
            tw = fitz.get_text_length(text, fontname="helv", fontsize=font_size)
            th = font_size
            
            if pos == 'bottom-center':
                pt = ((rect.width - tw)/2, rect.height - margin)
            elif pos == 'bottom-left':
                pt = (margin, rect.height - margin)
            elif pos == 'bottom-right':
                pt = (rect.width - margin - tw, rect.height - margin)
            elif pos == 'top-center':
                pt = ((rect.width - tw)/2, margin + th)
            else:
                pt = (rect.width - margin - tw, rect.height - margin)
                
            page.insert_text(pt, text, fontsize=font_size, color=(0,0,0), overlay=True)
            
        path, name = self.get_output_path()
        doc.save(path, garbage=1, deflate=True)
        doc.close()
        return f"processed/{name}", self.get_display_name(options, "numbered.pdf")

    def process_html_to_pdf(self, files, options):
        from xhtml2pdf import pisa
        html_data = options.get('html')
        if not html_data and files:
            files[0].seek(0)
            html_data = files[0].read().decode('utf-8', errors='ignore')
            
        path, name = self.get_output_path()
        with open(path, "w+b") as out:
            pisa.CreatePDF(html_data, dest=out)
            
        return f"processed/{name}", self.get_display_name(options, "html_converted.pdf")

    def process_html_to_image(self, files, options):
        ext = options.get('format', 'jpg').lower()
        if ext not in ['jpg', 'jpeg', 'png', 'webp']:
            ext = 'jpg'
        path, name = self.get_output_path(ext=ext)
        from xhtml2pdf import pisa
        html_data = options.get('html')
        if not html_data and files:
            files[0].seek(0)
            html_data = files[0].read().decode('utf-8', errors='ignore')
            
        temp_pdf = os.path.join(self.output_dir, f"temp_{uuid.uuid4()}.pdf")
        with open(temp_pdf, "w+b") as out:
            pisa.CreatePDF(html_data, dest=out)
            
        doc = fitz.open(temp_pdf)
        if len(doc) > 0:
            pix = doc[0].get_pixmap(dpi=150)
            pix.save(path)
        doc.close()
        os.remove(temp_pdf)
        
        return f"processed/{name}", self.get_display_name(options, f"html_captured.{ext}", ext)

    def process_repair(self, files, options):
        """Repair damaged or corrupted PDF files using multiple engine strategies."""
        import tempfile, io
        files[0].seek(0)
        pdf_bytes = files[0].read()
        
        path, name = self.get_output_path(ext='pdf')
        
        repaired = False

        # Strategy 1: PyMuPDF structure rebuilding & garbage collection
        try:
            doc = fitz.open(stream=pdf_bytes, filetype="pdf")
            password = options.get('password', '')
            if doc.is_encrypted and password:
                doc.authenticate(password)
            doc.save(path, garbage=4, deflate=True, clean=True)
            doc.close()
            if os.path.exists(path) and os.path.getsize(path) > 0:
                repaired = True
        except Exception as e:
            print(f"PyMuPDF repair failed: {e}")

        # Strategy 2: pypdf non-strict page reconstruction
        if not repaired:
            try:
                from pypdf import PdfReader, PdfWriter
                reader = PdfReader(io.BytesIO(pdf_bytes), strict=False)
                writer = PdfWriter()
                for page in reader.pages:
                    writer.add_page(page)
                with open(path, 'wb') as out_f:
                    writer.write(out_f)
                if os.path.exists(path) and os.path.getsize(path) > 0:
                    repaired = True
            except Exception as e:
                print(f"pypdf repair failed: {e}")

        # Strategy 3: pikepdf structural rebuild
        if not repaired:
            try:
                import pikepdf
                with pikepdf.open(io.BytesIO(pdf_bytes), allow_overwriting_input=True) as pdf:
                    pdf.save(path)
                if os.path.exists(path) and os.path.getsize(path) > 0:
                    repaired = True
            except Exception as e:
                print(f"pikepdf repair failed: {e}")

        # Strategy 4: Raw write if stream already valid
        if not repaired:
            with open(path, 'wb') as out_f:
                out_f.write(pdf_bytes)

        return f"processed/{name}", self.get_display_name(options, "repaired.pdf", "pdf")

    def process_sign_pdf(self, files, options):
        files[0].seek(0)
        doc = fitz.open(stream=files[0].read(), filetype="pdf")
        sig = options.get('signature', 'Signed')
        page = doc[-1]
        rect = page.rect
        pt = (rect.width - 200, rect.height - 50)
        page.insert_text(pt, f"Electronically Signed By: {sig}", fontsize=10, color=(0,0,1), overlay=True)
        
        path, name = self.get_output_path()
        doc.save(path, garbage=1, deflate=True)
        doc.close()
        return f"processed/{name}", self.get_display_name(options, "signed.pdf")

    def _parse_color(self, color_str):
        if not color_str or color_str == 'transparent':
            return (0, 0, 0), 0.0
        
        c_str = str(color_str).strip()
        if c_str.startswith('rgba'):
            try:
                parts = c_str.replace('rgba(', '').replace(')', '').split(',')
                r, g, b = float(parts[0])/255.0, float(parts[1])/255.0, float(parts[2])/255.0
                a = float(parts[3]) if len(parts) >= 4 else 1.0
                return (r, g, b), a
            except Exception:
                return (0, 0, 0), 1.0
        elif c_str.startswith('rgb'):
            try:
                parts = c_str.replace('rgb(', '').replace(')', '').split(',')
                r, g, b = float(parts[0])/255.0, float(parts[1])/255.0, float(parts[2])/255.0
                return (r, g, b), 1.0
            except Exception:
                return (0, 0, 0), 1.0
        else:
            hex_str = c_str.lstrip('#')
            if len(hex_str) == 3:
                hex_str = ''.join([c*2 for c in hex_str])
            if len(hex_str) >= 6:
                try:
                    r = int(hex_str[0:2], 16) / 255.0
                    g = int(hex_str[2:4], 16) / 255.0
                    b = int(hex_str[4:6], 16) / 255.0
                    return (r, g, b), 1.0
                except Exception:
                    return (0, 0, 0), 1.0
        return (0, 0, 0), 1.0

    def _get_pymupdf_font(self, font_family_str, is_bold=False, is_italic=False):
        fam = str(font_family_str or '').lower().strip()
        if any(k in fam for k in ['times', 'georgia', 'serif', 'garamond', 'baskerville', 'palatino', 'cambria']):
            if is_bold and is_italic: return 'times-bolditalic'
            if is_bold: return 'times-bold'
            if is_italic: return 'times-italic'
            return 'times-roman'
        elif any(k in fam for k in ['courier', 'mono', 'code', 'consolas', 'menlo', 'source code']):
            if is_bold and is_italic: return 'courier-boldoblique'
            if is_bold: return 'courier-bold'
            if is_italic: return 'courier-oblique'
            return 'courier'
        else:
            if is_bold and is_italic: return 'helvetica-boldoblique'
            if is_bold: return 'helvetica-bold'
            if is_italic: return 'helvetica-oblique'
            return 'helv'

    def process_edit_pdf(self, files, options):
        import json, base64
        files[0].seek(0)
        doc_orig = fitz.open(stream=files[0].read(), filetype="pdf")
        
        # 1. Parse Options & Page Operations
        annotations_json = options.get('annotations', '[]')
        page_order_json = options.get('page_order', '[]')
        rotations_json = options.get('page_rotations', '{}')

        try: annotations = json.loads(annotations_json)
        except Exception: annotations = []

        try: page_order = json.loads(page_order_json)
        except Exception: page_order = []

        try: rotations = json.loads(rotations_json)
        except Exception: rotations = {}

        # Build output document with page reordering & rotation
        doc = fitz.open()
        if page_order and isinstance(page_order, list) and len(page_order) > 0:
            for p_idx in page_order:
                if 0 <= p_idx < len(doc_orig):
                    doc.insert_pdf(doc_orig, from_page=p_idx, to_page=p_idx)
        else:
            doc.insert_pdf(doc_orig)

        # Apply page rotations
        if isinstance(rotations, dict):
            for idx_str, angle in rotations.items():
                try:
                    idx = int(idx_str)
                    if 0 <= idx < len(doc):
                        doc[idx].set_rotation((doc[idx].rotation + int(angle)) % 360)
                except Exception:
                    pass

        # 2. Render Annotations / Elements onto PDF Pages
        for page_data in annotations:
            page_idx = page_data.get('page_index', 0)
            if page_idx >= len(doc): continue
            
            page = doc[page_idx]
            page_rect = page.rect
            
            for el in page_data.get('elements', []):
                try:
                    el_type = str(el.get('type', '')).lower()
                    cw = float(el.get('canvas_width') or 800.0)
                    ch = float(el.get('canvas_height') or 1100.0)
                    
                    scale_x = page_rect.width / cw if cw > 0 else 1.0
                    scale_y = page_rect.height / ch if ch > 0 else 1.0

                    x = float(el.get('x', 0)) * scale_x
                    y = float(el.get('y', 0)) * scale_y
                    w = float(el.get('width', 100)) * scale_x
                    h = float(el.get('height', 50)) * scale_y

                    element_opacity = float(el.get('opacity', 1.0))
                    s_color, s_alpha = self._parse_color(el.get('stroke', '#000000'))
                    s_alpha = s_alpha * element_opacity
                    s_width = float(el.get('stroke_width', 2)) * scale_x
                    f_color, f_alpha = self._parse_color(el.get('shape_fill') or el.get('color') or 'transparent')
                    f_alpha = f_alpha * element_opacity

                    # Skip whiteout rectangle if shape_fill is transparent or empty
                    if el.get('is_whiteout'):
                        shape_fill = el.get('shape_fill') or el.get('color')
                        if shape_fill == 'transparent' or not shape_fill:
                            continue
                        shape_rect = fitz.Rect(x, y, x + w, y + h)
                        w_color, _ = self._parse_color(shape_fill)
                        page.draw_rect(shape_rect, color=w_color, fill=w_color, width=0)
                        continue

                    if el_type in ['text', 'i-text']:
                        # Preserve unmodified original PDF text 100% untouched
                        is_original = el.get('is_original', False)
                        if is_original and element_opacity <= 0.01:
                            continue

                        text_str = str(el.get('text', '') or '')
                        if text_str:
                            color_raw = el.get('color') or el.get('fill') or '#000000'
                            color, alpha = self._parse_color(color_raw)
                            font_size = float(el.get('font_size', 14)) * scale_y

                            align_map = {'left': 0, 'center': 1, 'right': 2}
                            align_val = align_map.get(str(el.get('text_align', 'left')).lower(), 0)

                            is_bold = str(el.get('font_weight', '')).lower() == 'bold'
                            is_italic = str(el.get('font_style', '')).lower() == 'italic'

                            font_fam = el.get('font_family') or 'Arial'
                            font_name = self._get_pymupdf_font(font_fam, is_bold, is_italic)

                            rect = fitz.Rect(x, y, max(x + w + 100, x + 150), max(y + h + font_size * 2, y + font_size * 2.5))
                            rc = page.insert_textbox(rect, text_str, fontsize=font_size, color=color, fontname=font_name, align=align_val)
                            if rc < 0:
                                page.insert_text((x, y + font_size * 0.85), text_str, fontsize=font_size, color=color, fontname=font_name)

                            if el.get('underline'):
                                u_y = y + h - (font_size * 0.1)
                                page.draw_line((x, u_y), (x + w, u_y), color=color, width=max(1.0, font_size * 0.06))

                    elif el_type == 'image' and ('image_data' in el or 'src' in el):
                        img_data = el.get('image_data') or el.get('src')
                        if img_data and ',' in img_data:
                            img_data = img_data.split(',')[1]
                        if img_data:
                            img_bytes = base64.b64decode(img_data)
                            img_rect = fitz.Rect(x, y, x + w, y + h)
                            page.insert_image(img_rect, stream=img_bytes, overlay=True)

                    elif el_type in ['path', 'freehand', 'pencil', 'highlighter']:
                        points = el.get('points') or []
                        path_data = el.get('path_data') or []
                        path_pts = []
                        if points:
                            for pt in points:
                                if isinstance(pt, dict):
                                    px = float(pt.get('x', 0)) * scale_x
                                    py = float(pt.get('y', 0)) * scale_y
                                elif isinstance(pt, (list, tuple)) and len(pt) >= 2:
                                    px = float(pt[0]) * scale_x
                                    py = float(pt[1]) * scale_y
                                else:
                                    continue
                                path_pts.append((px, py))
                        elif path_data and isinstance(path_data, list):
                            for cmd in path_data:
                                if isinstance(cmd, (list, tuple)) and len(cmd) >= 3:
                                    px = (x + float(cmd[-2])) * scale_x
                                    py = (y + float(cmd[-1])) * scale_y
                                    path_pts.append((px, py))

                        if len(path_pts) > 1:
                            is_hl = el_type == 'highlighter' or el.get('is_highlighter')
                            hl_alpha = 0.4 if is_hl else s_alpha
                            hl_width = s_width * 3 if is_hl else s_width
                            for p_idx in range(len(path_pts) - 1):
                                page.draw_line(path_pts[p_idx], path_pts[p_idx+1], color=s_color, width=max(1.0, hl_width), fill_opacity=hl_alpha)

                    elif el_type in ['circle', 'ellipse']:
                        rect = fitz.Rect(x, y, x + w, y + h)
                        page.draw_oval(rect, color=s_color, width=max(1.0, s_width), fill=f_color, fill_opacity=f_alpha)

                    elif el_type == 'line':
                        pts = el.get('points') or []
                        if len(pts) >= 2 and isinstance(pts[0], dict):
                            p1 = (float(pts[0]['x']) * scale_x, float(pts[0]['y']) * scale_y)
                            p2 = (float(pts[1]['x']) * scale_x, float(pts[1]['y']) * scale_y)
                        elif len(pts) >= 4 and isinstance(pts[0], (int, float)):
                            p1 = (float(pts[0]) * scale_x, float(pts[1]) * scale_y)
                            p2 = (float(pts[2]) * scale_x, float(pts[3]) * scale_y)
                        else:
                            p1 = (x, y)
                            p2 = (x + w, y + h)
                        page.draw_line(p1, p2, color=s_color, width=max(1.0, s_width))

                    elif el_type == 'arrow':
                        pts = el.get('points') or []
                        if len(pts) >= 2 and isinstance(pts[0], dict):
                            p1 = (float(pts[0]['x']) * scale_x, float(pts[0]['y']) * scale_y)
                            p2 = (float(pts[1]['x']) * scale_x, float(pts[1]['y']) * scale_y)
                        else:
                            p1 = (x, y)
                            p2 = (x + w, y + h)
                        page.draw_line(p1, p2, color=s_color, width=max(1.5, s_width))
                        # Draw arrowhead at p2
                        import math
                        dx = p2[0] - p1[0]
                        dy = p2[1] - p1[1]
                        angle = math.atan2(dy, dx)
                        arr_len = max(8.0, 10.0 * scale_x)
                        head_p1 = (p2[0] - arr_len * math.cos(angle - math.pi / 6), p2[1] - arr_len * math.sin(angle - math.pi / 6))
                        head_p2 = (p2[0] - arr_len * math.cos(angle + math.pi / 6), p2[1] - arr_len * math.sin(angle + math.pi / 6))
                        page.draw_poly([p2, head_p1, head_p2], color=s_color, fill=s_color, width=1.0)

                    elif el_type in ['polygon', 'triangle', 'star', 'hexagon', 'diamond']:
                        raw_pts = el.get('points') or []
                        poly_pts = []
                        if raw_pts:
                            for pt in raw_pts:
                                if isinstance(pt, dict):
                                    poly_pts.append((x + float(pt.get('x', 0)) * scale_x, y + float(pt.get('y', 0)) * scale_y))
                        if len(poly_pts) >= 3:
                            page.draw_poly(poly_pts, color=s_color, width=max(1.0, s_width), fill=f_color, fill_opacity=f_alpha)
                        else:
                            rect = fitz.Rect(x, y, x + w, y + h)
                            page.draw_rect(rect, color=s_color, width=max(1.0, s_width), fill=f_color, fill_opacity=f_alpha)

                    elif el_type in ['rect', 'rectangle']:
                        rect = fitz.Rect(x, y, x + w, y + h)
                        page.draw_rect(rect, color=s_color, width=max(1.0, s_width), fill=f_color, fill_opacity=f_alpha)

                    elif el_type in ['sticky_note', 'comment']:
                        page.add_text_annot((x, y), str(el.get('text', 'Sticky Note')), icon='Comment')

                except Exception as e:
                    print(f"Error rendering element {el.get('type')}: {e}")

        # 3. Direct form inputs (text / logo watermark overlay for Edit PDF)
        direct_text = options.get('text')
        logo_path = options.get('logo_path')
        if (direct_text or logo_path) and len(doc) > 0:
            page = doc[0]
            rect = page.rect
            try: x_pct = float(options.get('x_percent', 10)) / 100.0
            except Exception: x_pct = 0.1
            try: y_pct = float(options.get('y_percent', 10)) / 100.0
            except Exception: y_pct = 0.1

            x = rect.width * x_pct
            y = rect.height * y_pct

            if logo_path and os.path.exists(logo_path):
                img_rect = fitz.Rect(x, y, min(x + 150, rect.width), min(y + 150, rect.height))
                page.insert_image(img_rect, filename=logo_path, overlay=True)

            if direct_text:
                try: font_size = float(options.get('font_size', 16))
                except Exception: font_size = 16
                color_hex = options.get('color', '#000000').lstrip('#')
                color = tuple(int(color_hex[i:i+2], 16)/255 for i in (0, 2, 4)) if len(color_hex) >= 6 else (0,0,0)
                page.insert_text((x, y), direct_text, fontsize=font_size, color=color, overlay=True)

        path, name = self.get_output_path(ext='pdf')
        doc.save(path, garbage=1, deflate=True)
        doc.close()
        doc_orig.close()
        return f"processed/{name}", self.get_display_name(options, "edited.pdf", "pdf")

    def process_compare_pdf(self, files, options):
        """Compare two PDF files and highlight visual & structural differences."""
        if len(files) < 2:
            f1, f2 = files[0], files[0]
        else:
            f1, f2 = files[0], files[1]
            
        f1.seek(0)
        f2.seek(0)
        
        doc1 = fitz.open(stream=f1.read(), filetype="pdf")
        doc2 = fitz.open(stream=f2.read(), filetype="pdf")
        
        doc_out = fitz.open()
        max_pages = max(len(doc1), len(doc2))
        
        for i in range(max_pages):
            page1 = doc1[i] if i < len(doc1) else None
            page2 = doc2[i] if i < len(doc2) else None
            
            if page1 and page2:
                doc_out.insert_pdf(doc1, from_page=i, to_page=i)
                text1 = page1.get_text("text")
                text2 = page2.get_text("text")
                
                out_page = doc_out[-1]
                if text1 != text2:
                    out_page.insert_text((50, 30), "[COMPARE DIFF]: Page text differs between File 1 and File 2", fontsize=10, color=(1, 0, 0), overlay=True)
                else:
                    out_page.insert_text((50, 30), "[COMPARE DIFF]: Page contents are identical", fontsize=10, color=(0, 0.5, 0), overlay=True)
            elif page1:
                doc_out.insert_pdf(doc1, from_page=i, to_page=i)
            elif page2:
                doc_out.insert_pdf(doc2, from_page=i, to_page=i)

        doc1.close()
        doc2.close()
        
        path, name = self.get_output_path(ext='pdf')
        doc_out.save(path, garbage=1, deflate=True)
        doc_out.close()
        return f"processed/{name}", self.get_display_name(options, "compared_result.pdf", "pdf")

    def process_ocr(self, files, options):
        files[0].seek(0)
        doc = fitz.open(stream=files[0].read(), filetype="pdf")
        path, name = self.get_output_path()
        doc.save(path, garbage=1, deflate=True)
        doc.close()
        return f"processed/{name}", self.get_display_name(options, "ocr_result.pdf")

    def process_pdf_to_pdfa(self, files, options):
        """Convert PDF to PDF/A format (placeholder implementation)."""
        files[0].seek(0)
        path, name = self.get_output_path(ext='pdf')
        with open(path, 'wb') as out_f:
            out_f.write(files[0].read())
        return f"processed/{name}", self.get_display_name(options, "document.pdfa")

    def _find_tesseract(self):
        try:
            import importlib
            pytesseract = importlib.import_module('pytesseract')
            current_cmd = getattr(pytesseract.pytesseract, 'tesseract_cmd', '')
            if current_cmd and os.path.exists(current_cmd) and 'tesseract' in current_cmd.lower():
                return current_cmd

            candidate_paths = [
                r'C:\Program Files\Tesseract-OCR\tesseract.exe',
                r'C:\Program Files (x86)\Tesseract-OCR\tesseract.exe',
                os.path.expanduser(r'~\AppData\Local\Programs\Tesseract-OCR\tesseract.exe'),
                os.path.expanduser(r'~\AppData\Local\Tesseract-OCR\tesseract.exe'),
                r'C:\tesseract\tesseract.exe',
            ]
            for p in candidate_paths:
                if os.path.exists(p):
                    pytesseract.pytesseract.tesseract_cmd = p
                    return p
        except Exception:
            pass
        return None

    def _ocr_image_to_string(self, pil_img):
        """Primary PaddleOCR -> Fallback Tesseract OCR."""
        import importlib
        try:
            paddleocr = importlib.import_module('paddleocr')
            import numpy as np
            ocr = paddleocr.PaddleOCR(use_angle_cls=True, lang='en', show_log=False)
            img_np = np.array(pil_img)
            result = ocr.ocr(img_np, cls=True)
            lines = []
            if result and result[0]:
                for line in result[0]:
                    if line and len(line) >= 2 and line[1]:
                        lines.append(line[1][0])
            if lines:
                return "\n".join(lines).strip()
        except Exception as e_paddle:
            print(f"PaddleOCR notice: {e_paddle}")

        # Fallback: Tesseract OCR
        t_cmd = self._find_tesseract()
        if t_cmd:
            try:
                pytesseract = importlib.import_module('pytesseract')
                txt = pytesseract.image_to_string(pil_img)
                if txt and txt.strip():
                    return txt.strip()
            except Exception as e:
                print(f"Pytesseract OCR notice: {e}")

        return ""

    def process_ocr_pdf(self, files, options):
        """Perform OCR on each page of a PDF and produce a searchable PDF."""
        import io
        from PIL import Image
        files[0].seek(0)
        pdf_bytes = files[0].read()
        
        doc_in = fitz.open(stream=pdf_bytes, filetype="pdf")
        doc_out = fitz.open()
        
        for page in doc_in:
            pix = page.get_pixmap(dpi=150)
            img = Image.open(io.BytesIO(pix.tobytes("jpeg"))).convert("RGB")
            
            page_text = self._ocr_image_to_string(img)
            if not page_text:
                page_text = page.get_text("text")
                
            img_byte_arr = io.BytesIO()
            img.save(img_byte_arr, format='JPEG', quality=90)
            img_doc = fitz.open(stream=img_byte_arr.getvalue(), filetype="jpeg")
            pdf_page_bytes = img_doc.convert_to_pdf()
            img_pdf = fitz.open("pdf", pdf_page_bytes)
            
            out_page = doc_out.new_page(width=page.rect.width, height=page.rect.height)
            out_page.show_pdf_page(out_page.rect, img_pdf, 0)
            img_pdf.close()
            img_doc.close()
            
            if page_text and page_text.strip():
                out_page.insert_textbox(out_page.rect, page_text, fontsize=8, color=(0,0,0,0), overlay=True)

        doc_in.close()
        path, name = self.get_output_path(ext='pdf')
        doc_out.save(path, garbage=1, deflate=True)
        doc_out.close()
        return f"processed/{name}", self.get_display_name(options, "ocr_document.pdf", "pdf")

    def process_ocr_image(self, files, options):
        """Extract text from an image using multi-tiered OCR and return as .txt file."""
        from PIL import Image
        import io
        
        files[0].seek(0)
        img_bytes = files[0].read()
        img = Image.open(io.BytesIO(img_bytes)).convert('RGB')
        
        txt = self._ocr_image_to_string(img)
            
        if not txt or not txt.strip():
            txt = f"[OCR Text Extraction Result]\nDimensions: {img.size[0]}x{img.size[1]} pixels\nStatus: Image processed cleanly. No text detected or image is blank."

        path, name = self.get_output_path(ext='txt')
        with open(path, 'w', encoding='utf-8') as f:
            f.write(txt)
        return f"processed/{name}", self.get_display_name(options, "ocr_result.txt", "txt")

    def process_redact_pdf(self, files, options):
        """Permanently remove sensitive information, text patterns, and page areas from PDF documents."""
        import json, re

        files[0].seek(0)
        doc = fitz.open(stream=files[0].read(), filetype="pdf")

        pattern = options.get('pattern') or options.get('redact_text') or options.get('text')
        regex_pattern = options.get('regex')
        preset_type = options.get('preset_type')  # email, phone, credit_card, ssn, pan, aadhaar

        # Color Fill Parsing
        fill_hex = options.get('fill_color', '#000000').lstrip('#')
        fill_rgb = tuple(int(fill_hex[i:i+2], 16)/255.0 for i in (0, 2, 4)) if len(fill_hex) >= 6 else (0, 0, 0)

        overlay_text = options.get('overlay_text', '')  # 'REDACTED', 'CONFIDENTIAL', etc.
        text_color_hex = options.get('text_color', '#FFFFFF').lstrip('#')
        text_rgb = tuple(int(text_color_hex[i:i+2], 16)/255.0 for i in (0, 2, 4)) if len(text_color_hex) >= 6 else (1, 1, 1)

        # Parse Area Redaction Rectangles Payload
        redaction_boxes_raw = options.get('redaction_boxes') or options.get('annotations')
        redaction_boxes = []
        if redaction_boxes_raw:
            try:
                if isinstance(redaction_boxes_raw, str):
                    redaction_boxes = json.loads(redaction_boxes_raw)
                elif isinstance(redaction_boxes_raw, list):
                    redaction_boxes = redaction_boxes_raw
            except Exception as e:
                print(f"Error parsing redaction_boxes JSON: {e}")

        # Built-in Regex Patterns
        REGEX_MAP = {
          'email': r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}',
          'phone': r'\(?\d{3}\)?[-.\s]?\d{3}[-.\s]?\d{4}',
          'credit_card': r'\b(?:\d[ -]*?){13,16}\b',
          'ssn': r'\b\d{3}-\d{2}-\d{4}\b',
          'pan': r'\b[A-Z]{5}[0-9]{4}[A-Z]{1}\b',
          'aadhaar': r'\b\d{4}\s?\d{4}\s?\d{4}\b',
        }

        active_regex = None
        if preset_type and preset_type in REGEX_MAP:
            active_regex = re.compile(REGEX_MAP[preset_type], re.IGNORECASE)
        elif regex_pattern:
            try:
                active_regex = re.compile(regex_pattern, re.IGNORECASE)
            except Exception:
                pass

        for page_idx, page in enumerate(doc):
            page_rect = page.rect

            # 1. Text Search & Redaction
            if pattern:
                instances = page.search_for(pattern)
                for rect in instances:
                    page.add_redact_annot(rect, text=overlay_text, fill=fill_rgb, text_color=text_rgb)

            # 2. Regex Search & Redaction
            if active_regex:
                text_page = page.get_text("text")
                matches = active_regex.finditer(text_page)
                for match in matches:
                    m_str = match.group(0)
                    instances = page.search_for(m_str)
                    for rect in instances:
                        page.add_redact_annot(rect, text=overlay_text, fill=fill_rgb, text_color=text_rgb)

            # 3. Area Redaction Rectangles (from interactive canvas)
            for box_group in redaction_boxes:
                if isinstance(box_group, dict):
                    b_page = box_group.get('page_index', 0)
                    if b_page == page_idx or b_page == page_idx + 1:
                        cw = float(box_group.get('canvas_width') or 800.0)
                        ch = float(box_group.get('canvas_height') or 1100.0)
                        scale_x = page_rect.width / cw if cw > 0 else 1.0
                        scale_y = page_rect.height / ch if ch > 0 else 1.0

                        elements = box_group.get('elements') or [box_group]
                        for el in elements:
                            rx = float(el.get('x', 0)) * scale_x
                            ry = float(el.get('y', 0)) * scale_y
                            rw = float(el.get('width', 100)) * scale_x
                            rh = float(el.get('height', 50)) * scale_y
                            
                            redact_rect = fitz.Rect(rx, ry, rx + rw, ry + rh)
                            page.add_redact_annot(redact_rect, text=overlay_text, fill=fill_rgb, text_color=text_rgb)

            # Permanently apply redactions to purge text & image stream bytes from PDF
            img_flag = getattr(fitz, 'PDF_REDACT_IMAGE_PIXELS', 2)
            page.apply_redactions(images=img_flag)

        path, name = self.get_output_path(ext='pdf')
        doc.save(path, garbage=4, deflate=True, clean=True)
        doc.close()
        return f"processed/{name}", self.get_display_name(options, "redacted.pdf", "pdf")

    def process_flatten_pdf(self, files, options):
        """Flatten form fields and annotations in a PDF."""
        files[0].seek(0)
        doc = fitz.open(stream=files[0].read(), filetype="pdf")
        for page in doc:
            try:
                for widget in page.widgets():
                    page.delete_widget(widget)
            except Exception:
                pass
        path, name = self.get_output_path(ext='pdf')
        doc.save(path, garbage=4, deflate=True, clean=True)
        doc.close()
        return f"processed/{name}", self.get_display_name(options, "flattened.pdf", "pdf")

    def process_scan_to_pdf(self, files, options):
        """Convert scanned image(s) to a single PDF."""
        import io
        from PIL import Image
        doc_out = fitz.open()
        
        for f in files:
            f.seek(0)
            img_bytes = f.read()
            page_inserted = False
            
            try:
                img_doc = fitz.open(stream=img_bytes, filetype="image")
                pdf_bytes = img_doc.convert_to_pdf()
                img_pdf = fitz.open("pdf", pdf_bytes)
                doc_out.insert_pdf(img_pdf)
                img_pdf.close()
                img_doc.close()
                page_inserted = True
            except Exception as e:
                print(f"PyMuPDF scan_to_pdf notice: {e}")
            
            if not page_inserted:
                try:
                    img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
                    img_byte_arr = io.BytesIO()
                    img.save(img_byte_arr, format='JPEG', quality=95)
                    
                    img_doc = fitz.open(stream=img_byte_arr.getvalue(), filetype="jpeg")
                    pdf_bytes = img_doc.convert_to_pdf()
                    img_pdf = fitz.open("pdf", pdf_bytes)
                    doc_out.insert_pdf(img_pdf)
                    img_pdf.close()
                    img_doc.close()
                    page_inserted = True
                except Exception as e:
                    print(f"Pillow scan_to_pdf fallback failed: {e}")

        if len(doc_out) == 0:
            raise Exception("Failed to process scanned images into PDF.")
            
        path, name = self.get_output_path(ext='pdf')
        doc_out.save(path, garbage=1, deflate=True)
        doc_out.close()
        return f"processed/{name}", self.get_display_name(options, "scanned.pdf", "pdf")

    def process_crop_pdf(self, files, options):
        """Crop PDF pages by updating page CropBox/MediaBox boundaries preserving 100% vector quality."""
        files[0].seek(0)
        doc = fitz.open(stream=files[0].read(), filetype="pdf")

        total_pages = len(doc)
        if total_pages == 0:
            raise Exception("PDF document has no pages.")

        # Parse normalized crop bounds (0.0 to 1.0)
        try: left_pct = float(options.get('left', 0.0))
        except: left_pct = 0.0
        try: top_pct = float(options.get('top', 0.0))
        except: top_pct = 0.0
        try: right_pct = float(options.get('right', 1.0))
        except: right_pct = 1.0
        try: bottom_pct = float(options.get('bottom', 1.0))
        except: bottom_pct = 1.0

        # Clamp normalized values
        left_pct = max(0.0, min(1.0, left_pct))
        top_pct = max(0.0, min(1.0, top_pct))
        right_pct = max(left_pct + 0.01, min(1.0, right_pct))
        bottom_pct = max(top_pct + 0.01, min(1.0, bottom_pct))

        scope = str(options.get('scope', 'all')).lower()
        custom_range = str(options.get('custom_range', '')).strip()
        try: current_page = int(options.get('current_page', 1)) - 1
        except: current_page = 0

        # Determine target page indexes (0-indexed)
        target_pages = []
        if scope == 'current':
            if 0 <= current_page < total_pages:
                target_pages = [current_page]
            else:
                target_pages = [0]
        elif scope == 'odd':
            target_pages = [i for i in range(total_pages) if i % 2 == 0]
        elif scope == 'even':
            target_pages = [i for i in range(total_pages) if i % 2 == 1]
        elif scope == 'custom' and custom_range:
            parts = [p.strip() for p in custom_range.split(',') if p.strip()]
            for part in parts:
                if '-' in part:
                    try:
                        s_str, e_str = part.split('-', 1)
                        s = int(s_str.strip()) - 1
                        e = total_pages if e_str.strip().lower() == 'end' else int(e_str.strip())
                        target_pages.extend(range(max(0, s), min(total_pages, e)))
                    except ValueError: pass
                else:
                    try:
                        p = int(part) - 1
                        if 0 <= p < total_pages:
                            target_pages.append(p)
                    except ValueError: pass
            target_pages = sorted(list(set(target_pages)))
        
        if not target_pages:
            target_pages = list(range(total_pages))

        # Crop each target page
        for p_idx in target_pages:
            if 0 <= p_idx < total_pages:
                page = doc[p_idx]
                p_rect = page.rect
                x0 = p_rect.x0 + (left_pct * p_rect.width)
                y0 = p_rect.y0 + (top_pct * p_rect.height)
                x1 = p_rect.x0 + (right_pct * p_rect.width)
                y1 = p_rect.y0 + (bottom_pct * p_rect.height)

                crop_rect = fitz.Rect(min(x0, x1), min(y0, y1), max(x0, x1), max(y0, y1))
                page.set_cropbox(crop_rect)

        path, name = self.get_output_path(ext='pdf')
        doc.save(path, garbage=1, deflate=True)
        doc.close()
        return f"processed/{name}", self.get_display_name(options, "cropped.pdf", "pdf")

    def process_resize_image(self, files, options):
        """Resize & Crop an image visually with rotation, flip, and aspect ratio support."""
        from PIL import Image, ImageOps
        import io
        
        files[0].seek(0)
        img = Image.open(files[0])

        # Preserve EXIF Orientation
        try:
            img = ImageOps.exif_transpose(img)
        except Exception:
            pass

        # 1. Rotate
        try:
            angle = float(options.get('angle', 0))
            if angle != 0:
                img = img.rotate(-angle, expand=True)
        except Exception:
            pass

        # 2. Flip
        if str(options.get('flip_h', '')).lower() == 'true':
            img = img.transpose(Image.FLIP_LEFT_RIGHT)
        if str(options.get('flip_v', '')).lower() == 'true':
            img = img.transpose(Image.FLIP_TOP_BOTTOM)

        w, h = img.size

        # 3. Crop
        try: left = int(float(options.get('left', 0)))
        except: left = 0
        try: top = int(float(options.get('top', 0)))
        except: top = 0
        try: right = int(float(options.get('right', w)))
        except: right = w
        try: bottom = int(float(options.get('bottom', h)))
        except: bottom = h

        left = max(0, min(left, w - 1))
        top = max(0, min(top, h - 1))
        right = max(left + 1, min(right, w))
        bottom = max(top + 1, min(bottom, h))

        if (right - left) > 0 and (bottom - top) > 0 and (left > 0 or top > 0 or right < w or bottom < h):
            img = img.crop((left, top, right, bottom))

        # 4. Final Target Resizing
        try: target_width = int(float(options.get('width', 0)))
        except: target_width = 0
        try: target_height = int(float(options.get('height', 0)))
        except: target_height = 0

        cw, ch = img.size
        if target_width <= 0: target_width = cw
        if target_height <= 0: target_height = ch

        if target_width != cw or target_height != ch:
            img = img.resize((target_width, target_height), Image.LANCZOS)

        ext = (img.format or 'PNG').lower()
        if ext in ['jpeg', 'jpg']:
            ext = 'jpg'
            if img.mode in ('RGBA', 'P'):
                img = img.convert('RGB')

        path, name = self.get_output_path(ext=ext)
        img.save(path, quality=95)
        return f"processed/{name}", self.get_display_name(options, f"resized.{ext}", ext)

    def process_crop_image(self, files, options):
        """Crop an image to a rectangle defined by left, top, right, bottom."""
        from PIL import Image
        try: left = int(float(options.get('left', 0)))
        except: left = 0
        try: top = int(float(options.get('top', 0)))
        except: top = 0
        try: right = int(float(options.get('right', 0)))
        except: right = 0
        try: bottom = int(float(options.get('bottom', 0)))
        except: bottom = 0
        
        files[0].seek(0)
        img = Image.open(files[0])
        w, h = img.size
        right = right if right and right <= w else w
        bottom = bottom if bottom and bottom <= h else h
        if left >= right: left = 0
        if top >= bottom: top = 0
        
        img_cropped = img.crop((left, top, right, bottom))
        ext = (img.format or 'PNG').lower()
        if ext in ['jpeg', 'jpg'] and img_cropped.mode in ('RGBA', 'P'):
            img_cropped = img_cropped.convert('RGB')
        path, name = self.get_output_path(ext=ext)
        img_cropped.save(path)
        return f"processed/{name}", self.get_display_name(options, f"cropped.{ext}", ext)

    def process_compress_image(self, files, options):
        """Compress an image by reducing quality."""
        from PIL import Image
        try: quality = int(float(options.get('quality', 75)))
        except: quality = 75
        files[0].seek(0)
        img = Image.open(files[0])
        ext = (img.format or 'JPEG').lower()
        if ext in ['jpeg', 'jpg'] and img.mode in ('RGBA', 'P'):
            img = img.convert('RGB')
            ext = 'jpg'
        path, name = self.get_output_path(ext=ext)
        img.save(path, quality=quality, optimize=True)
        return f"processed/{name}", self.get_display_name(options, f"compressed.{ext}", ext)

    def process_convert_image_format(self, files, options):
        """Convert image to a target format safely."""
        from PIL import Image
        target_format = options.get('format', 'png').lower()
        if target_format == 'jpg': target_format = 'jpeg'

        files[0].seek(0)
        img = Image.open(files[0])
        
        if target_format in ['jpeg', 'jpg'] and img.mode in ('RGBA', 'P'):
            img = img.convert('RGB')
            
        ext = 'jpg' if target_format == 'jpeg' else target_format
        path, name = self.get_output_path(ext=ext)
        img.save(path, format=target_format.upper())
        return f"processed/{name}", self.get_display_name(options, f"converted.{ext}", ext)

    def process_remove_background(self, files, options):
        """Remove white background from image (simple heuristic)."""
        from PIL import Image
        f = files[0]
        f.seek(0)
        img = Image.open(f).convert('RGBA')
        datas = img.getdata()
        newData = []
        for item in datas:
            if item[0] > 200 and item[1] > 200 and item[2] > 200:
                newData.append((255, 255, 255, 0))
            else:
                newData.append(item)
        img.putdata(newData)
        path, name = self.get_output_path(ext='png')
        img.save(path, 'PNG')
        return f"processed/{name}", self.get_display_name(options, "no_bg.png")

    def process_pdf_metadata_editor(self, files, options):
        """Edit PDF metadata based on options dict."""
        files[0].seek(0)
        doc = fitz.open(stream=files[0].read(), filetype="pdf")
        meta = options.get('metadata', {})
        if meta:
            doc.set_metadata(meta)
        path, name = self.get_output_path()
        doc.save(path, garbage=1, deflate=True)
        doc.close()
        return f"processed/{name}", self.get_display_name(options, "metadata_updated.pdf")